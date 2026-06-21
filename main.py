#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
================================================================================
 PyPDF Studio  —  A professional desktop PDF application (single-file)
================================================================================

Inspired by Adobe Acrobat / Foxit / Nitro. Built with PySide6 (Qt6) and a stack
of PDF tooling (PyMuPDF, pypdf, Pillow, ReportLab, pytesseract, OpenCV).

Everything lives in this one file by design. It is organized into clearly
separated sections:

    1.  Imports & optional-dependency detection
    2.  Logging
    3.  Application constants & theming (QSS)
    4.  SettingsManager  (SQLite: settings, recent files, favorites, history)
    5.  PdfDocument      (thin wrapper around fitz.Document)
    6.  Background workers (QThreadPool) for heavy operations
    7.  Viewer widgets   (lazy-rendering page view, thumbnails)
    8.  Side panels      (thumbnails, bookmarks, comments, search, properties)
    9.  Tool dialogs     (merge, split, watermark, security, compress, OCR, ...)
   10.  MainWindow       (menus, toolbar, tabs, status bar, drag & drop)
   11.  Entry point

The application degrades gracefully: optional features (OCR, Office export,
OpenCV) are detected at runtime and disabled with friendly messages when their
backing library is not installed.

Required:    PySide6, PyMuPDF (fitz)
Recommended: pypdf, Pillow, reportlab
Optional:    pytesseract (+ Tesseract binary), opencv-python,
             python-docx, openpyxl, python-pptx

Run:         python main.py
================================================================================
"""

from __future__ import annotations

# ----------------------------------------------------------------------------
# 1. Imports & optional-dependency detection
# ----------------------------------------------------------------------------
import os
import io
import sys
import time
import json
import math
import shutil
import sqlite3
import logging
import tempfile
import datetime
import traceback
from dataclasses import dataclass, field
from typing import Optional, List, Dict, Tuple, Any, Callable

# --- Hard requirement: PySide6 ----------------------------------------------
try:
    from PySide6.QtCore import (
        Qt, QSize, QRect, QRectF, QPoint, QPointF, QTimer, QThreadPool,
        QRunnable, QObject, Signal, Slot, QSettings, QStandardPaths, QMimeData,
        QEvent, QBuffer, QByteArray,
    )
    from PySide6.QtGui import (
        QAction, QIcon, QPixmap, QImage, QPainter, QColor, QPen, QBrush, QFont,
        QKeySequence, QPalette, QCursor, QPolygonF, QFontDatabase, QPdfWriter,
        QPageSize, QPageLayout, QGuiApplication, QTransform, QActionGroup,
        QDragEnterEvent, QDropEvent, QShortcut, QTextDocument,
    )
    from PySide6.QtWidgets import (
        QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, QGridLayout,
        QLabel, QPushButton, QToolButton, QScrollArea, QTabWidget, QSplitter,
        QFileDialog, QMessageBox, QStatusBar, QToolBar, QMenu, QMenuBar,
        QListWidget, QListWidgetItem, QTreeWidget, QTreeWidgetItem, QLineEdit,
        QComboBox, QSpinBox, QDoubleSpinBox, QCheckBox, QRadioButton, QDialog,
        QDialogButtonBox, QFormLayout, QGroupBox, QSlider, QColorDialog,
        QProgressDialog, QProgressBar, QInputDialog, QStackedWidget, QFrame,
        QSizePolicy, QStyle, QStyleFactory, QPlainTextEdit, QTextEdit, QFontComboBox,
    )
    from PySide6.QtPrintSupport import QPrinter, QPrintDialog, QPrintPreviewDialog
except Exception as exc:  # pragma: no cover - fatal
    sys.stderr.write(
        "FATAL: PySide6 is required.\n"
        "Install it with:  pip install PySide6\n"
        f"Underlying error: {exc}\n"
    )
    raise

# --- Hard requirement: PyMuPDF (fitz) ---------------------------------------
try:
    import fitz  # PyMuPDF
    HAS_FITZ = True
except Exception:
    fitz = None
    HAS_FITZ = False

# --- Optional libraries ------------------------------------------------------
def _try_import(name: str):
    try:
        return __import__(name), True
    except Exception:
        return None, False


pypdf, HAS_PYPDF = _try_import("pypdf")
PIL, HAS_PIL = _try_import("PIL")
if HAS_PIL:
    try:
        from PIL import Image, ImageDraw, ImageFont, ImageQt  # noqa
    except Exception:
        HAS_PIL = False

reportlab, HAS_REPORTLAB = _try_import("reportlab")
pytesseract, HAS_TESSERACT = _try_import("pytesseract")
cv2, HAS_CV2 = _try_import("cv2")
docx, HAS_DOCX = _try_import("docx")
openpyxl, HAS_OPENPYXL = _try_import("openpyxl")
pptx, HAS_PPTX = _try_import("pptx")


# ----------------------------------------------------------------------------
# 2. Logging
# ----------------------------------------------------------------------------
def _app_base_dir() -> str:
    """Folder containing the running program (the .exe when frozen, else this
    script). This is the anchor for portable data storage."""
    if getattr(sys, "frozen", False):          # bundled by PyInstaller
        return os.path.dirname(sys.executable)
    return os.path.dirname(os.path.abspath(__file__))


def _app_data_dir() -> str:
    """
    Return (and create) a writable directory for app data.

    PORTABLE BEHAVIOUR: prefer a 'PyPDFStudio-Data' folder right next to the
    program, so the whole app (exe + settings + logs) can live on a USB stick
    and leave no traces on the host. If that location is not writable (e.g. the
    app was placed under Program Files), gracefully fall back to %APPDATA% and
    finally the temp directory.
    """
    candidates = [
        os.path.join(_app_base_dir(), "PyPDFStudio-Data"),
        os.path.join(os.environ.get("APPDATA") or os.path.expanduser("~"), "PyPDFStudio"),
        os.path.join(tempfile.gettempdir(), "PyPDFStudio"),
    ]
    for path in candidates:
        try:
            os.makedirs(path, exist_ok=True)
            probe = os.path.join(path, ".write_test")
            with open(probe, "w") as fh:
                fh.write("ok")
            os.remove(probe)
            return path
        except Exception:
            continue
    return tempfile.gettempdir()


APP_DIR = _app_data_dir()
LOG_PATH = os.path.join(APP_DIR, "pypdf_studio.log")

# Build handlers defensively. When frozen with --windowed there is no console,
# so sys.stdout/stderr are None — a StreamHandler bound to them would silently
# break logging. Only attach a console handler when a real stream exists.
_handlers: List[logging.Handler] = []
try:
    _handlers.append(logging.FileHandler(LOG_PATH, encoding="utf-8"))
except Exception:
    pass
if sys.stdout is not None:
    _handlers.append(logging.StreamHandler(sys.stdout))

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
    handlers=_handlers or None,
)
log = logging.getLogger("PyPDFStudio")


def install_excepthook(app_window_getter: Callable[[], Optional[QWidget]] = lambda: None):
    """Route uncaught exceptions to the log and a friendly dialog."""
    def hook(exc_type, exc_value, exc_tb):
        msg = "".join(traceback.format_exception(exc_type, exc_value, exc_tb))
        log.error("Uncaught exception:\n%s", msg)
        try:
            parent = app_window_getter()
            QMessageBox.critical(
                parent, "Unexpected Error",
                "An unexpected error occurred.\n\n"
                f"{exc_type.__name__}: {exc_value}\n\n"
                f"Details were written to:\n{LOG_PATH}",
            )
        except Exception:
            pass
    sys.excepthook = hook


# ----------------------------------------------------------------------------
# 3. Application constants & theming
# ----------------------------------------------------------------------------
APP_NAME = "PyPDF Studio"
APP_VERSION = "1.0.0"
ORG_NAME = "PyPDFStudio"

DEFAULT_ACCENT = "#3b82f6"

# Reading-mode page tints (multiplied onto rendered pages).
READING_MODES = {
    "Normal": None,
    "Sepia": (244, 232, 205),
    "Eye Comfort": (235, 240, 232),
    "Night": (40, 42, 48),
}


def _hex_to_rgb(h: str) -> Tuple[int, int, int]:
    h = h.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    try:
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    except Exception:
        return 59, 130, 246


def build_stylesheet(theme: str, accent: str, font: str = "Segoe UI") -> str:
    """Return a Qt stylesheet for the given theme ('dark'/'light') and accent."""
    if theme == "dark":
        bg = "#14161b"        # window background (deepest)
        canvas = "#0f1115"    # document viewing canvas
        bg2 = "#1c1f26"       # chrome: menu/tool/status bars, panels
        bg3 = "#262a33"       # inputs, buttons, raised surfaces
        fg = "#e8eaed"
        fg_dim = "#9aa1ad"
        border = "#2b3039"
    else:  # light
        bg = "#eceff3"
        canvas = "#dfe3ea"
        bg2 = "#ffffff"
        bg3 = "#f3f5f8"
        fg = "#1a1c20"
        fg_dim = "#5f6571"
        border = "#dde1e8"

    ar, ag, ab = _hex_to_rgb(accent)
    soft = f"rgba({ar}, {ag}, {ab}, 0.16)"      # accent wash (hover/selected)
    softer = f"rgba({ar}, {ag}, {ab}, 0.09)"
    accent_hi = f"rgba({ar}, {ag}, {ab}, 0.55)"  # focus ring

    return f"""
    * {{
        font-family: "{font}", "Segoe UI", "Inter", "Helvetica Neue", Arial, sans-serif;
        font-size: 13px;
        color: {fg};
        outline: none;
    }}
    QMainWindow, QDialog {{ background-color: {bg}; }}
    QWidget {{ background-color: {bg}; }}
    QToolTip {{
        background-color: {bg3}; color: {fg};
        border: 1px solid {border}; border-radius: 6px; padding: 5px 8px;
    }}

    QMenuBar {{
        background-color: {bg2};
        border-bottom: 1px solid {border};
        padding: 3px 6px;
    }}
    QMenuBar::item {{ padding: 6px 12px; background: transparent; border-radius: 7px; }}
    QMenuBar::item:selected {{ background: {soft}; color: {fg}; }}
    QMenu {{
        background-color: {bg2};
        border: 1px solid {border};
        padding: 6px;
        border-radius: 10px;
    }}
    QMenu::item {{ padding: 7px 28px 7px 14px; border-radius: 7px; }}
    QMenu::item:selected {{ background-color: {accent}; color: white; }}
    QMenu::icon {{ padding-left: 8px; }}
    QMenu::separator {{ height: 1px; background: {border}; margin: 6px 10px; }}

    QToolBar {{
        background-color: {bg2};
        border: none;
        border-bottom: 1px solid {border};
        spacing: 3px;
        padding: 6px 10px;
    }}
    QToolBar QToolButton {{
        background: transparent;
        border: 1px solid transparent;
        border-radius: 8px;
        padding: 6px;
        color: {fg};
    }}
    QToolBar QToolButton:hover {{ background-color: {soft}; }}
    QToolBar QToolButton:pressed {{ background-color: {accent}; }}
    QToolBar QToolButton:checked {{ background-color: {accent}; color: white; }}
    QToolBar QToolButton:disabled {{ background: transparent; }}
    QToolBar::separator {{
        background: {border}; width: 1px; margin: 6px 8px;
    }}

    QStatusBar {{
        background-color: {bg2};
        border-top: 1px solid {border};
        color: {fg_dim};
        padding: 3px 8px;
    }}
    QStatusBar::item {{ border: none; }}
    QLabel#StatusChip {{
        background-color: {bg3};
        border: 1px solid {border};
        border-radius: 9px;
        padding: 2px 11px;
        color: {fg_dim};
        font-size: 12px;
    }}
    QLabel#StatusMsg {{ color: {fg_dim}; padding-left: 4px; }}

    QTabWidget::pane {{ border: none; background: {bg}; top: -1px; }}
    QTabBar {{ qproperty-drawBase: 0; background: {bg2}; }}
    QTabBar::tab {{
        background: transparent;
        border: 1px solid transparent;
        padding: 8px 18px;
        margin: 4px 2px 0 2px;
        border-radius: 9px;
        color: {fg_dim};
    }}
    QTabBar::tab:selected {{ background: {bg}; color: {fg}; border: 1px solid {border}; }}
    QTabBar::tab:hover:!selected {{ background: {soft}; color: {fg}; }}
    QTabBar::close-button {{
        subcontrol-position: right;
        border-radius: 5px;
    }}
    QTabBar::close-button:hover {{ background: {soft}; }}

    QSplitter::handle {{ background: {bg}; }}
    QSplitter::handle:horizontal {{ width: 4px; }}
    QSplitter::handle:vertical {{ height: 4px; }}
    QSplitter::handle:hover {{ background: {accent_hi}; }}

    QListWidget, QTreeWidget, QPlainTextEdit, QTextEdit {{
        background-color: {bg2};
        border: 1px solid {border};
        border-radius: 10px;
        padding: 3px;
    }}
    QScrollArea {{ background-color: {canvas}; border: none; }}
    QListWidget::item, QTreeWidget::item {{ padding: 6px 6px; border-radius: 7px; }}
    QListWidget::item:selected, QTreeWidget::item:selected {{
        background-color: {accent}; color: white;
    }}
    QListWidget::item:hover:!selected, QTreeWidget::item:hover:!selected {{
        background-color: {soft};
    }}

    QHeaderView::section {{
        background-color: {bg3}; color: {fg_dim};
        border: none; border-bottom: 1px solid {border}; padding: 6px;
    }}

    QLineEdit, QComboBox, QSpinBox, QDoubleSpinBox, QFontComboBox {{
        background-color: {bg3};
        border: 1px solid {border};
        border-radius: 8px;
        padding: 6px 9px;
        selection-background-color: {accent};
        selection-color: white;
    }}
    QLineEdit:focus, QComboBox:focus, QSpinBox:focus, QDoubleSpinBox:focus,
    QFontComboBox:focus, QPlainTextEdit:focus, QTextEdit:focus {{
        border: 1px solid {accent};
    }}
    QComboBox::drop-down {{ border: none; width: 22px; }}
    QComboBox QAbstractItemView {{
        background-color: {bg2}; border: 1px solid {border}; border-radius: 8px;
        padding: 4px; selection-background-color: {accent}; selection-color: white;
    }}

    QPushButton {{
        background-color: {bg3};
        border: 1px solid {border};
        border-radius: 8px;
        padding: 8px 18px;
        color: {fg};
    }}
    QPushButton:hover {{ background-color: {soft}; border-color: {accent_hi}; }}
    QPushButton:pressed {{ background-color: {accent}; color: white; }}
    QPushButton:disabled {{ color: {fg_dim}; background: {bg2}; border-color: {border}; }}
    QPushButton#accent {{
        background-color: {accent}; color: white; border: none; font-weight: 600;
        padding: 9px 22px;
    }}
    QPushButton#accent:hover {{ background-color: {accent}; border: none; }}

    QGroupBox {{
        border: 1px solid {border};
        border-radius: 10px;
        margin-top: 12px;
        padding: 12px 10px 10px 10px;
        font-weight: 600;
        background-color: {softer};
    }}
    QGroupBox::title {{
        subcontrol-origin: margin;
        left: 14px; padding: 0 6px;
        color: {fg_dim};
    }}

    QCheckBox, QRadioButton {{ spacing: 8px; }}
    QCheckBox::indicator, QRadioButton::indicator {{
        width: 17px; height: 17px;
        border: 1px solid {border}; background: {bg3};
    }}
    QCheckBox::indicator {{ border-radius: 5px; }}
    QRadioButton::indicator {{ border-radius: 9px; }}
    QCheckBox::indicator:checked, QRadioButton::indicator:checked {{
        background-color: {accent}; border-color: {accent};
    }}
    QCheckBox::indicator:hover, QRadioButton::indicator:hover {{ border-color: {accent}; }}

    QSlider::groove:horizontal {{ height: 5px; background: {border}; border-radius: 3px; }}
    QSlider::handle:horizontal {{
        background: {accent}; width: 16px; height: 16px;
        margin: -6px 0; border-radius: 8px; border: 2px solid {bg2};
    }}
    QSlider::sub-page:horizontal {{ background: {accent}; border-radius: 3px; }}

    QProgressBar {{
        border: none; border-radius: 8px; text-align: center;
        background: {bg3}; height: 16px; color: {fg};
    }}
    QProgressBar::chunk {{ background-color: {accent}; border-radius: 8px; }}

    QScrollBar:vertical {{ background: transparent; width: 12px; margin: 2px; }}
    QScrollBar::handle:vertical {{ background: {border}; border-radius: 5px; min-height: 32px; }}
    QScrollBar::handle:vertical:hover {{ background: {fg_dim}; }}
    QScrollBar:horizontal {{ background: transparent; height: 12px; margin: 2px; }}
    QScrollBar::handle:horizontal {{ background: {border}; border-radius: 5px; min-width: 32px; }}
    QScrollBar::handle:horizontal:hover {{ background: {fg_dim}; }}
    QScrollBar::add-line, QScrollBar::sub-line {{ height: 0; width: 0; }}
    QScrollBar::add-page, QScrollBar::sub-page {{ background: transparent; }}

    QLabel#PageCanvas {{ background: transparent; }}
    QFrame#SidebarRail {{ background-color: {bg2}; border-right: 1px solid {border}; }}
    QFrame#PanelHeader {{ border-bottom: 1px solid {border}; }}
    QLabel#PanelTitle {{ font-weight: 700; font-size: 13px; color: {fg}; padding: 8px 6px; }}
    QToolButton#RailButton {{
        padding: 9px 2px; border-radius: 10px; font-size: 10px;
        color: {fg_dim}; font-weight: 600;
    }}
    QToolButton#RailButton:hover {{ background-color: {soft}; color: {fg}; }}
    QToolButton#RailButton:checked {{ background-color: {accent}; color: white; }}

    /* Welcome screen */
    QFrame#WelcomeCard {{
        background-color: {bg2};
        border: 1px solid {border};
        border-radius: 16px;
    }}
    QLabel#WelcomeTitle {{ font-size: 30px; font-weight: 800; color: {fg}; }}
    QLabel#WelcomeSub {{ color: {fg_dim}; font-size: 14px; }}
    QLabel#SectionLabel {{ color: {fg_dim}; font-size: 12px; font-weight: 700; }}
    QListWidget#RecentList {{ background: transparent; border: none; }}
    QPushButton#TileButton {{
        text-align: left; padding: 16px; border-radius: 12px;
        background-color: {bg3}; border: 1px solid {border};
        font-size: 14px; font-weight: 600; color: {fg};
    }}
    QPushButton#TileButton:hover {{ border-color: {accent}; background-color: {soft}; }}
    """


# Icon system ---------------------------------------------------------------
# Crisp, line-based vector icons drawn with QPainter on a 24x24 logical grid.
# This keeps the app dependency-free (no icon asset files) while looking far
# cleaner than glyph/emoji fonts, which render inconsistently across systems.
ICON_COLOR = "#8b92a1"   # neutral mid-tone that reads on both dark & light


def vector_icon(name: str, color: str = ICON_COLOR, size: int = 44) -> QIcon:
    """Render a named line icon to a QIcon. Unknown names get a rounded square."""
    pm = QPixmap(size, size)
    pm.fill(Qt.transparent)
    p = QPainter(pm)
    try:
        p.setRenderHint(QPainter.Antialiasing, True)
        s = size / 24.0
        col = QColor(color)
        pen = QPen(col)
        pen.setWidthF(2.0 * s)
        pen.setCapStyle(Qt.RoundCap)
        pen.setJoinStyle(Qt.RoundJoin)
        p.setPen(pen)
        p.setBrush(Qt.NoBrush)

        def line(x1, y1, x2, y2):
            p.drawLine(QPointF(x1 * s, y1 * s), QPointF(x2 * s, y2 * s))

        def rrect(x, y, w, h, r=2.5):
            p.drawRoundedRect(QRectF(x * s, y * s, w * s, h * s), r * s, r * s)

        def ellipse(cx, cy, rx, ry=None):
            ry = rx if ry is None else ry
            p.drawEllipse(QPointF(cx * s, cy * s), rx * s, ry * s)

        def polyline(pts, close=False):
            poly = QPolygonF([QPointF(a * s, b * s) for a, b in pts])
            p.drawPolygon(poly) if close else p.drawPolyline(poly)

        def arc(x, y, w, h, start_deg, span_deg):
            p.drawArc(QRectF(x * s, y * s, w * s, h * s),
                      int(start_deg * 16), int(span_deg * 16))

        def fill_poly(pts):
            p.setBrush(QBrush(col))
            p.drawPolygon(QPolygonF([QPointF(a * s, b * s) for a, b in pts]))
            p.setBrush(Qt.NoBrush)

        if name == "new":
            polyline([(6, 3), (14, 3), (18, 7), (18, 21), (6, 21)], close=True)
            polyline([(14, 3), (14, 7), (18, 7)])
        elif name == "open":
            polyline([(3, 8), (10, 8), (12, 6), (20, 6)])
            polyline([(3, 8), (3, 19), (21, 19), (21, 10), (8, 10), (3, 19)])
        elif name == "save":
            polyline([(5, 4), (17, 4), (20, 7), (20, 20), (5, 20)], close=True)
            line(8, 4, 8, 9); line(16, 4, 16, 9)
            rrect(8, 13, 9, 7, 1)
        elif name == "print":
            polyline([(7, 3), (17, 3), (17, 8)])
            line(7, 8, 7, 3)
            rrect(4, 8, 16, 8, 2)
            rrect(7, 14, 10, 6, 1)
            ellipse(16.5, 11, 0.7)
        elif name == "export":
            rrect(4, 4, 16, 16, 3)
            line(12, 14, 12, 5)
            polyline([(9, 8), (12, 5), (15, 8)])
        elif name == "undo":
            arc(4, 5, 15, 14, 30, 230)
            polyline([(4, 6), (4, 11), (9, 11)])
        elif name == "redo":
            arc(5, 5, 15, 14, -80, 230)
            polyline([(20, 6), (20, 11), (15, 11)])
        elif name == "search":
            ellipse(10.5, 10.5, 6.5)
            line(15.5, 15.5, 20, 20)
        elif name == "zoom_in":
            ellipse(10, 10, 6.2)
            line(14.6, 14.6, 20, 20)
            line(10, 7, 10, 13); line(7, 10, 13, 10)
        elif name == "zoom_out":
            ellipse(10, 10, 6.2)
            line(14.6, 14.6, 20, 20)
            line(7, 10, 13, 10)
        elif name == "fit_width":
            line(4, 4, 4, 20); line(20, 4, 20, 20)
            line(8, 12, 16, 12)
            polyline([(10, 9), (7, 12), (10, 15)])
            polyline([(14, 9), (17, 12), (14, 15)])
        elif name == "fit_page":
            rrect(4, 4, 16, 16, 2)
            polyline([(8, 8), (11, 8)]); polyline([(8, 8), (8, 11)])
            polyline([(16, 16), (13, 16)]); polyline([(16, 16), (16, 13)])
        elif name == "rotate_left":
            arc(4, 4, 16, 16, 60, 250)
            polyline([(4, 5), (5, 10), (10, 9)])
        elif name == "rotate_right":
            arc(4, 4, 16, 16, -130, 250)
            polyline([(20, 5), (19, 10), (14, 9)])
        elif name == "fullscreen":
            polyline([(4, 8), (4, 4), (8, 4)])
            polyline([(16, 4), (20, 4), (20, 8)])
            polyline([(20, 16), (20, 20), (16, 20)])
            polyline([(8, 20), (4, 20), (4, 16)])
        elif name == "highlight":
            fill_poly([(5, 19), (5, 16), (14, 7), (17, 10), (8, 19)])
            line(4, 21, 20, 21)
        elif name == "underline":
            line(7, 4, 7, 12); line(17, 4, 17, 12)
            arc(7, 8, 10, 9, 180, 180)
            line(5, 20, 19, 20)
        elif name == "strikeout":
            line(8, 5, 16, 5); line(12, 5, 12, 13)
            line(4, 14, 20, 14)
        elif name == "note":
            polyline([(4, 4), (20, 4), (20, 14), (13, 14), (8, 19), (8, 14), (4, 14)],
                     close=True)
            line(8, 8, 16, 8); line(8, 11, 13, 11)
        elif name == "image":
            rrect(3, 5, 18, 14, 2.5)
            ellipse(8, 9.5, 1.6)
            polyline([(4, 17), (9, 12), (12, 15), (16, 11), (20, 16)])
        elif name == "text":
            line(5, 5, 19, 5); line(12, 5, 12, 19)
            line(9, 19, 15, 19)
        elif name in ("page", "pages"):
            rrect(8, 3, 11, 15, 2)
            polyline([(5, 6), (5, 21), (15, 21)])
        elif name == "bookmark":
            polyline([(7, 3), (17, 3), (17, 21), (12, 16.5), (7, 21)], close=True)
        elif name == "comment":
            rrect(3, 4, 18, 13, 4)
            polyline([(8, 17), (8, 21), (13, 17)])
            line(7, 9, 17, 9); line(7, 12.5, 13, 12.5)
        elif name == "theme":
            ellipse(12, 12, 7)
            arc(5, 5, 14, 14, -90, 180)
            p.setBrush(QBrush(col))
            path_pts = [(12, 5)]
            # filled half
            p.setBrush(Qt.NoBrush)
            line(12, 5, 12, 19)
        elif name == "settings":
            ellipse(12, 12, 3.2)
            for k in range(8):
                ang = math.radians(k * 45)
                x1 = 12 + 5.2 * math.cos(ang); y1 = 12 + 5.2 * math.sin(ang)
                x2 = 12 + 7.4 * math.cos(ang); y2 = 12 + 7.4 * math.sin(ang)
                line(x1, y1, x2, y2)
        elif name == "merge":
            polyline([(4, 5), (10, 5), (12, 12), (4, 12)])
            polyline([(4, 12), (10, 12), (12, 19), (4, 19)])
            polyline([(12, 12), (20, 12)])
            polyline([(17, 9), (20, 12), (17, 15)])
        elif name == "split":
            line(12, 4, 12, 20)
            polyline([(4, 8), (8, 8), (8, 16), (4, 16)])
            polyline([(20, 8), (16, 8), (16, 16), (20, 16)])
        elif name == "lock":
            rrect(5, 10, 14, 10, 2.5)
            arc(8, 4, 8, 10, 0, 180)
            ellipse(12, 14.5, 1.3)
        else:
            rrect(5, 5, 14, 14, 3)
    finally:
        p.end()
    return QIcon(pm)


# Backwards-compatible alias (older call sites passed a glyph string).
def text_icon(name: str, color: str = ICON_COLOR, size: int = 44) -> QIcon:
    return vector_icon(name, color, size)


# ----------------------------------------------------------------------------
# 4. SettingsManager (SQLite-backed)
# ----------------------------------------------------------------------------
class SettingsManager:
    """
    Persists application settings, recent files, favorites and an operation
    history in a small SQLite database. All access is wrapped in try/except so
    a corrupt DB never crashes the app.
    """

    DEFAULTS = {
        "theme": "dark",
        "accent": DEFAULT_ACCENT,
        "default_zoom": "100",
        "startup_page": "welcome",   # 'welcome' | 'blank' | 'last'
        "view_mode": "continuous",   # 'single' | 'continuous' | 'two'
        "reading_mode": "Normal",
        "autosave": "1",
        "autosave_interval": "180",  # seconds
        "ui_font": "Segoe UI",
        "language": "English",
        "show_toolbar": "1",
    }

    def __init__(self, db_path: str):
        self.db_path = db_path
        self._conn = None
        self._init_db()

    def _init_db(self):
        try:
            self._conn = sqlite3.connect(self.db_path, check_same_thread=False)
            cur = self._conn.cursor()
            cur.execute("CREATE TABLE IF NOT EXISTS settings (key TEXT PRIMARY KEY, value TEXT)")
            cur.execute(
                "CREATE TABLE IF NOT EXISTS recent ("
                "path TEXT PRIMARY KEY, opened_at REAL, favorite INTEGER DEFAULT 0)"
            )
            cur.execute(
                "CREATE TABLE IF NOT EXISTS history ("
                "id INTEGER PRIMARY KEY AUTOINCREMENT, action TEXT, "
                "detail TEXT, ts REAL)"
            )
            self._conn.commit()
        except Exception as e:
            log.error("Failed to init settings DB: %s", e)
            self._conn = None

    # -- settings key/value --------------------------------------------------
    def get(self, key: str, default: Any = None) -> Any:
        if default is None:
            default = self.DEFAULTS.get(key)
        if not self._conn:
            return default
        try:
            cur = self._conn.execute("SELECT value FROM settings WHERE key=?", (key,))
            row = cur.fetchone()
            return row[0] if row else default
        except Exception as e:
            log.error("settings get(%s): %s", key, e)
            return default

    def get_int(self, key: str, default: int = 0) -> int:
        try:
            return int(self.get(key, default))
        except Exception:
            return default

    def get_bool(self, key: str, default: bool = False) -> bool:
        return str(self.get(key, "1" if default else "0")) in ("1", "true", "True")

    def set(self, key: str, value: Any):
        if not self._conn:
            return
        try:
            self._conn.execute(
                "INSERT INTO settings(key,value) VALUES(?,?) "
                "ON CONFLICT(key) DO UPDATE SET value=excluded.value",
                (key, str(value)),
            )
            self._conn.commit()
        except Exception as e:
            log.error("settings set(%s): %s", key, e)

    # -- recent files --------------------------------------------------------
    def add_recent(self, path: str):
        if not self._conn or not path:
            return
        try:
            self._conn.execute(
                "INSERT INTO recent(path, opened_at) VALUES(?,?) "
                "ON CONFLICT(path) DO UPDATE SET opened_at=excluded.opened_at",
                (path, time.time()),
            )
            self._conn.commit()
        except Exception as e:
            log.error("add_recent: %s", e)

    def recent_files(self, limit: int = 15) -> List[str]:
        if not self._conn:
            return []
        try:
            cur = self._conn.execute(
                "SELECT path FROM recent ORDER BY opened_at DESC LIMIT ?", (limit,)
            )
            return [r[0] for r in cur.fetchall() if os.path.exists(r[0])]
        except Exception as e:
            log.error("recent_files: %s", e)
            return []

    def clear_recent(self):
        if self._conn:
            try:
                self._conn.execute("DELETE FROM recent WHERE favorite=0")
                self._conn.commit()
            except Exception as e:
                log.error("clear_recent: %s", e)

    def toggle_favorite(self, path: str):
        if not self._conn:
            return
        try:
            cur = self._conn.execute("SELECT favorite FROM recent WHERE path=?", (path,))
            row = cur.fetchone()
            new = 0 if (row and row[0]) else 1
            self._conn.execute(
                "INSERT INTO recent(path, opened_at, favorite) VALUES(?,?,?) "
                "ON CONFLICT(path) DO UPDATE SET favorite=?",
                (path, time.time(), new, new),
            )
            self._conn.commit()
        except Exception as e:
            log.error("toggle_favorite: %s", e)

    def favorites(self) -> List[str]:
        if not self._conn:
            return []
        try:
            cur = self._conn.execute(
                "SELECT path FROM recent WHERE favorite=1 ORDER BY opened_at DESC"
            )
            return [r[0] for r in cur.fetchall() if os.path.exists(r[0])]
        except Exception:
            return []

    # -- history -------------------------------------------------------------
    def log_action(self, action: str, detail: str = ""):
        if self._conn:
            try:
                self._conn.execute(
                    "INSERT INTO history(action, detail, ts) VALUES(?,?,?)",
                    (action, detail, time.time()),
                )
                self._conn.commit()
            except Exception:
                pass

    def history(self, limit: int = 100) -> List[Tuple[str, str, float]]:
        if not self._conn:
            return []
        try:
            cur = self._conn.execute(
                "SELECT action, detail, ts FROM history ORDER BY ts DESC LIMIT ?", (limit,)
            )
            return cur.fetchall()
        except Exception:
            return []


# ----------------------------------------------------------------------------
# 5. PdfDocument — thin wrapper around fitz.Document
# ----------------------------------------------------------------------------
class PdfDocument:
    """
    Wraps a fitz.Document and tracks per-tab UI state (zoom, rotation, mode,
    current page, dirty flag). All mutating operations set `dirty = True`.
    """

    def __init__(self, path: Optional[str] = None):
        self.path: Optional[str] = path
        self.doc: Optional["fitz.Document"] = None
        self.dirty: bool = False
        self.password: Optional[str] = None
        # Per-tab view state
        self.zoom: float = 1.0
        self.rotation: int = 0          # 0/90/180/270 applied at render time
        self.current_page: int = 0
        self.view_mode: str = "continuous"

        if path:
            self.open(path)
        else:
            self.new_blank()

    # -- lifecycle -----------------------------------------------------------
    def open(self, path: str, password: Optional[str] = None):
        self.doc = fitz.open(path)
        if self.doc.needs_pass:
            ok = self.doc.authenticate(password or "")
            if not ok:
                raise PermissionError("Password required or incorrect.")
            self.password = password
        self.path = path
        self.dirty = False

    def new_blank(self, width: float = 595, height: float = 842):
        """Create a new single-page (A4) document."""
        self.doc = fitz.open()
        self.doc.new_page(width=width, height=height)
        self.path = None
        self.dirty = True

    def close(self):
        try:
            if self.doc:
                self.doc.close()
        except Exception:
            pass
        self.doc = None

    # -- queries -------------------------------------------------------------
    @property
    def page_count(self) -> int:
        return self.doc.page_count if self.doc else 0

    @property
    def name(self) -> str:
        if self.path:
            return os.path.basename(self.path)
        return "Untitled.pdf"

    def page_size(self, index: int) -> Tuple[float, float]:
        try:
            r = self.doc[index].rect
            return r.width, r.height
        except Exception:
            return 595.0, 842.0

    def metadata(self) -> Dict[str, str]:
        try:
            return dict(self.doc.metadata or {})
        except Exception:
            return {}

    def toc(self) -> List[list]:
        try:
            return self.doc.get_toc(simple=True) or []
        except Exception:
            return []

    # -- rendering -----------------------------------------------------------
    def render_page(self, index: int, zoom: float, rotation: int = 0,
                    tint: Optional[Tuple[int, int, int]] = None) -> QImage:
        """Render a page to a QImage at the given zoom and rotation."""
        page = self.doc[index]
        mat = fitz.Matrix(zoom, zoom)
        if rotation:
            mat = mat * fitz.Matrix(rotation)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        fmt = QImage.Format_RGB888
        img = QImage(pix.samples, pix.width, pix.height, pix.stride, fmt).copy()
        if tint:
            img = self._apply_tint(img, tint)
        return img

    @staticmethod
    def _apply_tint(img: QImage, tint: Tuple[int, int, int]) -> QImage:
        """Multiply a tint color over the image (for reading modes)."""
        out = img.convertToFormat(QImage.Format_RGB32)
        painter = QPainter(out)
        painter.setCompositionMode(QPainter.CompositionMode_Multiply)
        painter.fillRect(out.rect(), QColor(*tint))
        painter.end()
        return out

    # -- saving --------------------------------------------------------------
    def save(self, path: Optional[str] = None, **kwargs):
        target = path or self.path
        if not target:
            raise ValueError("No path supplied for save.")
        incremental = (target == self.path) and not kwargs
        try:
            if incremental and self.path:
                self.doc.save(target, incremental=True, encryption=fitz.PDF_ENCRYPT_KEEP)
            else:
                self.doc.save(target, garbage=4, deflate=True, **kwargs)
        except Exception:
            # Fall back to a full save to a temp file then move.
            tmp = target + ".tmp"
            self.doc.save(tmp, garbage=4, deflate=True)
            self.doc.close()
            shutil.move(tmp, target)
            self.doc = fitz.open(target)
        self.path = target
        self.dirty = False


# ----------------------------------------------------------------------------
# 6. Background workers
# ----------------------------------------------------------------------------
class WorkerSignals(QObject):
    finished = Signal(object)        # result
    error = Signal(str)
    progress = Signal(int, str)      # percent, message


class Worker(QRunnable):
    """Runs a callable on the global QThreadPool with progress callbacks."""

    def __init__(self, fn: Callable, *args, **kwargs):
        super().__init__()
        self.fn = fn
        self.args = args
        self.kwargs = kwargs
        self.signals = WorkerSignals()

    @Slot()
    def run(self):
        try:
            # Provide an optional progress callback to the worker function.
            self.kwargs.setdefault("progress_cb", self._progress)
            result = self.fn(*self.args, **self.kwargs)
            self.signals.finished.emit(result)
        except Exception as e:
            log.error("Worker error: %s\n%s", e, traceback.format_exc())
            self.signals.error.emit(str(e))

    def _progress(self, percent: int, message: str = ""):
        self.signals.progress.emit(percent, message)


# ----------------------------------------------------------------------------
# 7. Viewer widgets
# ----------------------------------------------------------------------------
class PageLabel(QLabel):
    """A single rendered page; tracks its page index and rendered state."""

    def __init__(self, index: int):
        super().__init__()
        self.index = index
        self.rendered_zoom = -1.0
        self.setAlignment(Qt.AlignCenter)
        self.setObjectName("PageCanvas")
        self.setFrameShape(QFrame.Shape.NoFrame)


class PdfViewer(QScrollArea):
    """
    Lazy-rendering PDF viewer.

    - Holds one PageLabel per page in a vertical (or grid for two-page) layout.
    - Only pages near the viewport are actually rendered (lazy loading).
    - Supports zoom, rotation, fit-width/page, single/continuous/two-page modes,
      Ctrl+wheel zoom, and current-page tracking for the status bar.
    """

    page_changed = Signal(int)       # 0-based current page
    zoom_changed = Signal(float)

    def __init__(self, document: PdfDocument, settings: SettingsManager):
        super().__init__()
        self.document = document
        self.settings = settings
        self.reading_mode = settings.get("reading_mode", "Normal")
        self.setWidgetResizable(True)
        self.setAlignment(Qt.AlignCenter)
        self.setBackgroundRole(QPalette.Dark)

        self._container = QWidget()
        self._container.setStyleSheet("background: transparent;")
        self.setWidget(self._container)
        self._labels: List[PageLabel] = []
        self._layout_mode = document.view_mode

        # Debounced lazy render on scroll/resize.
        self._render_timer = QTimer(self)
        self._render_timer.setSingleShot(True)
        self._render_timer.setInterval(40)
        self._render_timer.timeout.connect(self._render_visible)
        self.verticalScrollBar().valueChanged.connect(self._on_scroll)

        # Auto-scroll (reading mode)
        self._autoscroll_timer = QTimer(self)
        self._autoscroll_timer.timeout.connect(self._auto_scroll_step)
        self._autoscroll_speed = 1

        self.rebuild()

    # -- layout construction -------------------------------------------------
    def rebuild(self):
        """Recreate page placeholders according to the current view mode."""
        old = self._container.layout()
        if old:
            QWidget().setLayout(old)  # detach old layout
        self._labels.clear()

        mode = self.document.view_mode
        self._layout_mode = mode

        if mode == "two":
            layout = QGridLayout()
            layout.setSpacing(16)
            layout.setContentsMargins(20, 20, 20, 20)
            row = 0
            col = 0
            for i in range(self.document.page_count):
                lbl = PageLabel(i)
                self._labels.append(lbl)
                layout.addWidget(lbl, row, col)
                col += 1
                if col >= 2:
                    col = 0
                    row += 1
        else:
            layout = QVBoxLayout()
            layout.setSpacing(16)
            layout.setContentsMargins(20, 20, 20, 20)
            layout.setAlignment(Qt.AlignHCenter | Qt.AlignTop)
            for i in range(self.document.page_count):
                lbl = PageLabel(i)
                self._labels.append(lbl)
                layout.addWidget(lbl)

        self._container.setLayout(layout)
        self._size_placeholders()
        QTimer.singleShot(0, self._render_visible)

    def _size_placeholders(self):
        """Set each label's fixed size based on page size * zoom (rotation aware)."""
        zoom = self.document.zoom
        rot = self.document.rotation
        for lbl in self._labels:
            w, h = self.document.page_size(lbl.index)
            if rot in (90, 270):
                w, h = h, w
            lbl.setFixedSize(int(w * zoom) + 2, int(h * zoom) + 2)
            lbl.rendered_zoom = -1.0  # force re-render
            # Placeholder background
            lbl.setStyleSheet(
                "background: white; border: 1px solid rgba(0,0,0,0.25);"
                "border-radius: 2px;"
            )

    # -- visibility & lazy rendering -----------------------------------------
    def _on_scroll(self, _value):
        self._render_timer.start()
        self._update_current_page()

    def resizeEvent(self, event):
        super().resizeEvent(event)
        self._render_timer.start()

    def _visible_indices(self) -> List[int]:
        vp = self.viewport().rect()
        vp_global = QRect(self.widget().mapFromParent(QPoint(0, 0)), vp.size())
        view_top = self.verticalScrollBar().value()
        view_bottom = view_top + self.viewport().height()
        margin = self.viewport().height()  # render one screen ahead/behind
        result = []
        for lbl in self._labels:
            top = lbl.y()
            bottom = top + lbl.height()
            if bottom >= view_top - margin and top <= view_bottom + margin:
                result.append(lbl.index)
        return result

    def _render_visible(self):
        if not self.document.doc:
            return
        zoom = self.document.zoom
        rot = self.document.rotation
        tint = READING_MODES.get(self.reading_mode)
        for idx in self._visible_indices():
            lbl = self._labels[idx]
            if abs(lbl.rendered_zoom - zoom) < 1e-3:
                continue
            try:
                img = self.document.render_page(idx, zoom, rot, tint)
                pix = QPixmap.fromImage(img)
                lbl.setPixmap(pix)
                lbl.setFixedSize(pix.size())
                lbl.rendered_zoom = zoom
                lbl.setStyleSheet("background: transparent;")
            except Exception as e:
                log.error("render page %d: %s", idx, e)

    def _update_current_page(self):
        center = self.verticalScrollBar().value() + self.viewport().height() // 2
        best = 0
        for lbl in self._labels:
            if lbl.y() <= center:
                best = lbl.index
        if best != self.document.current_page:
            self.document.current_page = best
            self.page_changed.emit(best)

    # -- navigation ----------------------------------------------------------
    def goto_page(self, index: int):
        index = max(0, min(index, self.document.page_count - 1))
        if index < len(self._labels):
            lbl = self._labels[index]
            self.verticalScrollBar().setValue(lbl.y() - 20)
            self.document.current_page = index
            self.page_changed.emit(index)
            self._render_timer.start()

    def next_page(self):
        self.goto_page(self.document.current_page + 1)

    def prev_page(self):
        self.goto_page(self.document.current_page - 1)

    # -- zoom ----------------------------------------------------------------
    def set_zoom(self, zoom: float, anchor: Optional[QPoint] = None):
        zoom = max(0.1, min(zoom, 8.0))
        self.document.zoom = zoom
        self._size_placeholders()
        self._render_timer.start()
        self.zoom_changed.emit(zoom)

    def zoom_in(self):
        self.set_zoom(self.document.zoom * 1.15)

    def zoom_out(self):
        self.set_zoom(self.document.zoom / 1.15)

    def fit_width(self):
        if not self.document.page_count:
            return
        w, h = self.document.page_size(self.document.current_page)
        if self.document.rotation in (90, 270):
            w, h = h, w
        avail = self.viewport().width() - 60
        self.set_zoom(avail / w)

    def fit_page(self):
        if not self.document.page_count:
            return
        w, h = self.document.page_size(self.document.current_page)
        if self.document.rotation in (90, 270):
            w, h = h, w
        aw = self.viewport().width() - 60
        ah = self.viewport().height() - 60
        self.set_zoom(min(aw / w, ah / h))

    # -- rotation & modes ----------------------------------------------------
    def rotate(self, delta: int):
        self.document.rotation = (self.document.rotation + delta) % 360
        self._size_placeholders()
        self._render_timer.start()

    def set_view_mode(self, mode: str):
        self.document.view_mode = mode
        self.rebuild()

    def set_reading_mode(self, mode: str):
        self.reading_mode = mode
        for lbl in self._labels:
            lbl.rendered_zoom = -1.0
        self._render_timer.start()

    # -- auto scroll ---------------------------------------------------------
    def toggle_autoscroll(self, on: bool, speed: int = 1):
        self._autoscroll_speed = speed
        if on:
            self._autoscroll_timer.start(30)
        else:
            self._autoscroll_timer.stop()

    def _auto_scroll_step(self):
        sb = self.verticalScrollBar()
        if sb.value() >= sb.maximum():
            self._autoscroll_timer.stop()
            return
        sb.setValue(sb.value() + self._autoscroll_speed)

    # -- events --------------------------------------------------------------
    def wheelEvent(self, event):
        if event.modifiers() & Qt.ControlModifier:
            delta = event.angleDelta().y()
            if delta > 0:
                self.zoom_in()
            else:
                self.zoom_out()
            event.accept()
        else:
            super().wheelEvent(event)


# ----------------------------------------------------------------------------
# 8. Side panels
# ----------------------------------------------------------------------------
class ThumbnailPanel(QWidget):
    """Page thumbnails with drag-to-reorder and context actions."""

    page_selected = Signal(int)
    order_changed = Signal(list)     # new order of page indices

    def __init__(self):
        super().__init__()
        layout = QVBoxLayout(self)
        layout.setContentsMargins(4, 4, 4, 4)
        title = QLabel("Pages")
        title.setStyleSheet("font-weight: 600; padding: 4px;")
        layout.addWidget(title)

        self.list = QListWidget()
        self.list.setViewMode(QListWidget.IconMode)
        self.list.setIconSize(QSize(120, 160))
        self.list.setResizeMode(QListWidget.Adjust)
        self.list.setMovement(QListWidget.Snap)
        self.list.setDragDropMode(QListWidget.InternalMove)
        self.list.setSpacing(8)
        self.list.itemClicked.connect(self._on_click)
        self.list.model().rowsMoved.connect(self._on_moved)
        layout.addWidget(self.list)

    def populate(self, document: PdfDocument):
        self.list.clear()
        if not document.doc:
            return
        for i in range(document.page_count):
            try:
                img = document.render_page(i, 0.2, document.rotation)
                icon = QIcon(QPixmap.fromImage(img))
            except Exception:
                icon = QIcon()
            item = QListWidgetItem(icon, f"{i + 1}")
            item.setData(Qt.UserRole, i)
            item.setTextAlignment(Qt.AlignCenter)
            self.list.addItem(item)

    def _on_click(self, item):
        self.page_selected.emit(item.data(Qt.UserRole))

    def _on_moved(self, *args):
        order = [self.list.item(i).data(Qt.UserRole) for i in range(self.list.count())]
        self.order_changed.emit(order)


class BookmarkPanel(QWidget):
    """Outline / table of contents view."""

    bookmark_selected = Signal(int)   # 0-based page

    def __init__(self):
        super().__init__()
        layout = QVBoxLayout(self)
        layout.setContentsMargins(4, 4, 4, 4)
        title = QLabel("Bookmarks")
        title.setStyleSheet("font-weight: 600; padding: 4px;")
        layout.addWidget(title)
        self.tree = QTreeWidget()
        self.tree.setHeaderHidden(True)
        self.tree.itemClicked.connect(self._on_click)
        layout.addWidget(self.tree)

    def populate(self, document: PdfDocument):
        self.tree.clear()
        toc = document.toc()
        if not toc:
            placeholder = QTreeWidgetItem(["(No bookmarks)"])
            self.tree.addTopLevelItem(placeholder)
            return
        stack = [(0, None)]
        for level, title, page in toc:
            item = QTreeWidgetItem([title])
            item.setData(0, Qt.UserRole, max(0, page - 1))
            while stack and stack[-1][0] >= level:
                stack.pop()
            parent = stack[-1][1] if stack else None
            if parent is None:
                self.tree.addTopLevelItem(item)
            else:
                parent.addChild(item)
            stack.append((level, item))
        self.tree.expandAll()

    def _on_click(self, item, _col):
        page = item.data(0, Qt.UserRole)
        if page is not None:
            self.bookmark_selected.emit(page)


class CommentsPanel(QWidget):
    """Lists annotations found in the document."""

    annotation_selected = Signal(int)  # page

    def __init__(self):
        super().__init__()
        layout = QVBoxLayout(self)
        layout.setContentsMargins(4, 4, 4, 4)
        title = QLabel("Comments")
        title.setStyleSheet("font-weight: 600; padding: 4px;")
        layout.addWidget(title)
        self.list = QListWidget()
        self.list.itemClicked.connect(self._on_click)
        layout.addWidget(self.list)

    def populate(self, document: PdfDocument):
        self.list.clear()
        if not document.doc:
            return
        for pno in range(document.page_count):
            try:
                page = document.doc[pno]
                annot = page.first_annot
                while annot:
                    kind = annot.type[1] if annot.type else "Annot"
                    content = (annot.info.get("content") or "").strip()
                    label = f"Pg {pno + 1} · {kind}"
                    if content:
                        label += f": {content[:40]}"
                    item = QListWidgetItem(label)
                    item.setData(Qt.UserRole, pno)
                    self.list.addItem(item)
                    annot = annot.next
            except Exception:
                continue
        if self.list.count() == 0:
            self.list.addItem("(No comments)")

    def _on_click(self, item):
        page = item.data(Qt.UserRole)
        if page is not None:
            self.annotation_selected.emit(page)


class SearchPanel(QWidget):
    """Full-text search across the document with result navigation."""

    result_selected = Signal(int, object)   # page, fitz.Rect

    def __init__(self):
        super().__init__()
        self._document: Optional[PdfDocument] = None
        self._results: List[Tuple[int, Any]] = []
        layout = QVBoxLayout(self)
        layout.setContentsMargins(4, 4, 4, 4)
        title = QLabel("Search")
        title.setStyleSheet("font-weight: 600; padding: 4px;")
        layout.addWidget(title)

        row = QHBoxLayout()
        self.edit = QLineEdit()
        self.edit.setPlaceholderText("Find in document…")
        self.edit.returnPressed.connect(self.search)
        row.addWidget(self.edit)
        btn = QPushButton("Go")
        btn.clicked.connect(self.search)
        row.addWidget(btn)
        layout.addLayout(row)

        nav = QHBoxLayout()
        self.prev_btn = QToolButton(); self.prev_btn.setText("◀ Prev")
        self.next_btn = QToolButton(); self.next_btn.setText("Next ▶")
        self.prev_btn.clicked.connect(self.prev_result)
        self.next_btn.clicked.connect(self.next_result)
        nav.addWidget(self.prev_btn)
        nav.addWidget(self.next_btn)
        layout.addLayout(nav)

        self.count_label = QLabel("")
        self.count_label.setStyleSheet("color: gray; padding: 2px;")
        layout.addWidget(self.count_label)

        self.list = QListWidget()
        self.list.itemClicked.connect(self._on_item)
        layout.addWidget(self.list)
        self._cursor = -1

    def set_document(self, document: PdfDocument):
        self._document = document
        self._results.clear()
        self.list.clear()
        self.count_label.clear()

    def search(self):
        if not self._document or not self._document.doc:
            return
        query = self.edit.text().strip()
        if not query:
            return
        self._results.clear()
        self.list.clear()
        for pno in range(self._document.page_count):
            try:
                rects = self._document.doc[pno].search_for(query)
                for r in rects:
                    self._results.append((pno, r))
                    item = QListWidgetItem(f"Page {pno + 1}")
                    item.setData(Qt.UserRole, len(self._results) - 1)
                    self.list.addItem(item)
            except Exception:
                continue
        self.count_label.setText(f"{len(self._results)} match(es)")
        self._cursor = 0 if self._results else -1
        if self._results:
            self._emit(0)

    def _emit(self, i: int):
        if 0 <= i < len(self._results):
            pno, rect = self._results[i]
            self.result_selected.emit(pno, rect)
            self.list.setCurrentRow(i)

    def next_result(self):
        if self._results:
            self._cursor = (self._cursor + 1) % len(self._results)
            self._emit(self._cursor)

    def prev_result(self):
        if self._results:
            self._cursor = (self._cursor - 1) % len(self._results)
            self._emit(self._cursor)

    def _on_item(self, item):
        self._cursor = item.data(Qt.UserRole)
        self._emit(self._cursor)


class PropertiesPanel(QWidget):
    """Right sidebar: document metadata & quick info."""

    def __init__(self):
        super().__init__()
        layout = QVBoxLayout(self)
        layout.setContentsMargins(8, 8, 8, 8)
        title = QLabel("Properties")
        title.setStyleSheet("font-weight: 600; font-size: 14px; padding: 4px;")
        layout.addWidget(title)
        self.form = QFormLayout()
        self.form.setLabelAlignment(Qt.AlignRight)
        layout.addLayout(self.form)
        layout.addStretch(1)
        self._rows: Dict[str, QLabel] = {}

    def update_for(self, document: Optional[PdfDocument]):
        while self.form.rowCount():
            self.form.removeRow(0)
        if not document or not document.doc:
            self.form.addRow(QLabel("No document open."))
            return
        meta = document.metadata()
        size = "-"
        if document.path and os.path.exists(document.path):
            size = human_size(os.path.getsize(document.path))
        info = [
            ("File", document.name),
            ("Pages", str(document.page_count)),
            ("Size", size),
            ("Title", meta.get("title") or "-"),
            ("Author", meta.get("author") or "-"),
            ("Subject", meta.get("subject") or "-"),
            ("Creator", meta.get("creator") or "-"),
            ("Producer", meta.get("producer") or "-"),
            ("Created", meta.get("creationDate") or "-"),
            ("Encrypted", "Yes" if document.doc.is_encrypted else "No"),
        ]
        for k, v in info:
            val = QLabel(str(v))
            val.setWordWrap(True)
            val.setTextInteractionByMouse = True
            self.form.addRow(QLabel(f"<b>{k}</b>"), val)


# ----------------------------------------------------------------------------
# Utility helpers
# ----------------------------------------------------------------------------
def human_size(num: float) -> str:
    for unit in ("B", "KB", "MB", "GB", "TB"):
        if abs(num) < 1024.0:
            return f"{num:3.1f} {unit}"
        num /= 1024.0
    return f"{num:.1f} PB"


def parse_page_ranges(text: str, max_page: int) -> List[int]:
    """Parse '1-3,5,7-9' (1-based) into a sorted list of 0-based indices."""
    pages = set()
    text = (text or "").strip()
    if not text:
        return list(range(max_page))
    for part in text.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, _, b = part.partition("-")
            try:
                start = int(a)
                end = int(b)
                for p in range(start, end + 1):
                    if 1 <= p <= max_page:
                        pages.add(p - 1)
            except ValueError:
                continue
        else:
            try:
                p = int(part)
                if 1 <= p <= max_page:
                    pages.add(p - 1)
            except ValueError:
                continue
    return sorted(pages)


# ----------------------------------------------------------------------------
# 9. Tool dialogs
# ----------------------------------------------------------------------------
class WatermarkDialog(QDialog):
    """Configure and apply a text or image watermark to a document."""

    def __init__(self, document: PdfDocument, parent=None):
        super().__init__(parent)
        self.document = document
        self.setWindowTitle("Watermark")
        self.setMinimumWidth(420)
        form = QFormLayout(self)

        self.type_combo = QComboBox()
        self.type_combo.addItems(["Text", "Image"])
        self.type_combo.currentTextChanged.connect(self._toggle)
        form.addRow("Type", self.type_combo)

        self.text_edit = QLineEdit("CONFIDENTIAL")
        form.addRow("Text", self.text_edit)

        self.image_row = QHBoxLayout()
        self.image_path = QLineEdit()
        browse = QPushButton("Browse…")
        browse.clicked.connect(self._browse_image)
        self.image_row.addWidget(self.image_path)
        self.image_row.addWidget(browse)
        self._image_widget = QWidget(); self._image_widget.setLayout(self.image_row)
        form.addRow("Image", self._image_widget)

        self.opacity = QSlider(Qt.Horizontal)
        self.opacity.setRange(5, 100); self.opacity.setValue(25)
        form.addRow("Opacity %", self.opacity)

        self.rotation = QSpinBox(); self.rotation.setRange(-180, 180); self.rotation.setValue(45)
        form.addRow("Rotation", self.rotation)

        self.scale = QDoubleSpinBox(); self.scale.setRange(0.1, 5.0); self.scale.setValue(1.0)
        self.scale.setSingleStep(0.1)
        form.addRow("Scale", self.scale)

        self.position = QComboBox()
        self.position.addItems(["Center", "Top", "Bottom", "Tile"])
        form.addRow("Position", self.position)

        self.fontsize = QSpinBox(); self.fontsize.setRange(8, 200); self.fontsize.setValue(48)
        form.addRow("Font size", self.fontsize)

        self.color_btn = QPushButton("Pick color")
        self._color = QColor("#888888")
        self.color_btn.clicked.connect(self._pick_color)
        form.addRow("Color", self.color_btn)

        self.pages_edit = QLineEdit()
        self.pages_edit.setPlaceholderText("All pages (e.g. 1-3,5)")
        form.addRow("Pages", self.pages_edit)

        bb = QDialogButtonBox(QDialogButtonBox.Ok | QDialogButtonBox.Cancel)
        bb.accepted.connect(self.accept)
        bb.rejected.connect(self.reject)
        form.addRow(bb)
        self._toggle("Text")

    def _toggle(self, t):
        self.text_edit.setEnabled(t == "Text")
        self._image_widget.setEnabled(t == "Image")

    def _browse_image(self):
        path, _ = QFileDialog.getOpenFileName(self, "Image", "", "Images (*.png *.jpg *.jpeg)")
        if path:
            self.image_path.setText(path)

    def _pick_color(self):
        c = QColorDialog.getColor(self._color, self, "Watermark color")
        if c.isValid():
            self._color = c

    def apply(self):
        pages = parse_page_ranges(self.pages_edit.text(), self.document.page_count)
        opacity = self.opacity.value() / 100.0
        rot = self.rotation.value()
        pos = self.position.currentText()
        is_text = self.type_combo.currentText() == "Text"
        color = (self._color.redF(), self._color.greenF(), self._color.blueF())

        for pno in pages:
            page = self.document.doc[pno]
            rect = page.rect
            if is_text:
                text = self.text_edit.text()
                fs = self.fontsize.value() * self.scale.value()
                if pos == "Tile":
                    step = fs * 5
                    y = rect.height * 0.15
                    while y < rect.height:
                        x = rect.width * 0.1
                        while x < rect.width:
                            self._stamp_text(page, text, fitz.Point(x, y), fs, rot, color, opacity)
                            x += step
                        y += step
                else:
                    if pos == "Top":
                        pt = fitz.Point(rect.width / 2, rect.height * 0.15)
                    elif pos == "Bottom":
                        pt = fitz.Point(rect.width / 2, rect.height * 0.9)
                    else:
                        pt = fitz.Point(rect.width / 2, rect.height / 2)
                    self._stamp_text(page, text, pt, fs, rot, color, opacity, center=True)
            else:
                img_path = self.image_path.text()
                if img_path and os.path.exists(img_path):
                    w = rect.width * 0.5 * self.scale.value()
                    h = w
                    cx, cy = rect.width / 2, rect.height / 2
                    r = fitz.Rect(cx - w / 2, cy - h / 2, cx + w / 2, cy + h / 2)
                    page.insert_image(r, filename=img_path, overlay=True, keep_proportion=True)
        self.document.dirty = True

    @staticmethod
    def _stamp_text(page, text, point, fontsize, rotate, color, opacity, center=False):
        try:
            if center:
                tw = fitz.get_text_length(text, fontsize=fontsize)
                point = fitz.Point(point.x - tw / 2, point.y)
            page.insert_text(
                point, text, fontsize=fontsize, rotate=rotate,
                color=color, fill_opacity=opacity, overlay=True,
            )
        except Exception as e:
            log.error("watermark stamp: %s", e)


class SecurityDialog(QDialog):
    """Encrypt / decrypt / set permissions on the document."""

    def __init__(self, document: PdfDocument, parent=None):
        super().__init__(parent)
        self.document = document
        self.setWindowTitle("Security")
        self.setMinimumWidth(400)
        form = QFormLayout(self)

        self.user_pw = QLineEdit(); self.user_pw.setEchoMode(QLineEdit.Password)
        form.addRow("User password", self.user_pw)
        self.owner_pw = QLineEdit(); self.owner_pw.setEchoMode(QLineEdit.Password)
        form.addRow("Owner password", self.owner_pw)

        self.no_print = QCheckBox("Disable printing")
        self.no_copy = QCheckBox("Disable copying")
        self.no_edit = QCheckBox("Disable editing")
        form.addRow(self.no_print)
        form.addRow(self.no_copy)
        form.addRow(self.no_edit)

        note = QLabel("Leave passwords empty and click OK while choosing\n"
                      "“Remove security” to strip encryption.")
        note.setStyleSheet("color: gray;")
        form.addRow(note)

        self.remove = QCheckBox("Remove security instead")
        form.addRow(self.remove)

        bb = QDialogButtonBox(QDialogButtonBox.Ok | QDialogButtonBox.Cancel)
        bb.accepted.connect(self.accept)
        bb.rejected.connect(self.reject)
        form.addRow(bb)

    def build_save_kwargs(self) -> dict:
        if self.remove.isChecked():
            return dict(encryption=fitz.PDF_ENCRYPT_NONE)
        perm = int(
            fitz.PDF_PERM_ACCESSIBILITY
            | (0 if self.no_print.isChecked() else fitz.PDF_PERM_PRINT)
            | (0 if self.no_copy.isChecked() else fitz.PDF_PERM_COPY)
            | (0 if self.no_edit.isChecked() else fitz.PDF_PERM_MODIFY)
            | fitz.PDF_PERM_ANNOTATE
        )
        return dict(
            encryption=fitz.PDF_ENCRYPT_AES_256,
            owner_pw=self.owner_pw.text() or self.user_pw.text(),
            user_pw=self.user_pw.text(),
            permissions=perm,
        )


class HeaderFooterDialog(QDialog):
    """Add page numbers / date / time / filename / custom text."""

    def __init__(self, document: PdfDocument, parent=None):
        super().__init__(parent)
        self.document = document
        self.setWindowTitle("Header & Footer")
        form = QFormLayout(self)

        self.header_edit = QLineEdit()
        self.header_edit.setPlaceholderText("Header text (use {page} {pages} {date} {time} {file})")
        form.addRow("Header", self.header_edit)

        self.footer_edit = QLineEdit("Page {page} of {pages}")
        form.addRow("Footer", self.footer_edit)

        self.fontsize = QSpinBox(); self.fontsize.setRange(6, 48); self.fontsize.setValue(10)
        form.addRow("Font size", self.fontsize)

        bb = QDialogButtonBox(QDialogButtonBox.Ok | QDialogButtonBox.Cancel)
        bb.accepted.connect(self.accept)
        bb.rejected.connect(self.reject)
        form.addRow(bb)

    def apply(self):
        now = datetime.datetime.now()
        file = self.document.name
        total = self.document.page_count
        fs = self.fontsize.value()
        for pno in range(total):
            page = self.document.doc[pno]
            rect = page.rect
            ctx = dict(page=pno + 1, pages=total,
                       date=now.strftime("%Y-%m-%d"),
                       time=now.strftime("%H:%M"), file=file)
            header = self._fmt(self.header_edit.text(), ctx)
            footer = self._fmt(self.footer_edit.text(), ctx)
            if header:
                page.insert_text(fitz.Point(40, 30), header, fontsize=fs, overlay=True)
            if footer:
                page.insert_text(fitz.Point(40, rect.height - 25), footer, fontsize=fs, overlay=True)
        self.document.dirty = True

    @staticmethod
    def _fmt(template: str, ctx: dict) -> str:
        if not template:
            return ""
        try:
            return template.format(**ctx)
        except Exception:
            return template


class CompressDialog(QDialog):
    """Compress the PDF at a chosen level, showing size estimates."""

    LEVELS = {
        "Low": dict(dpi=150, quality=85),
        "Medium": dict(dpi=110, quality=70),
        "High": dict(dpi=90, quality=55),
        "Maximum": dict(dpi=72, quality=40),
    }

    def __init__(self, document: PdfDocument, parent=None):
        super().__init__(parent)
        self.document = document
        self.result_path: Optional[str] = None
        self.setWindowTitle("Compress PDF")
        self.setMinimumWidth(380)
        layout = QVBoxLayout(self)

        orig = "-"
        if document.path and os.path.exists(document.path):
            orig = human_size(os.path.getsize(document.path))
        layout.addWidget(QLabel(f"Original size: <b>{orig}</b>"))

        row = QHBoxLayout()
        row.addWidget(QLabel("Level"))
        self.level = QComboBox(); self.level.addItems(list(self.LEVELS.keys()))
        self.level.setCurrentText("Medium")
        row.addWidget(self.level)
        layout.addLayout(row)

        self.result_label = QLabel("")
        self.result_label.setStyleSheet("color: gray;")
        layout.addWidget(self.result_label)

        bb = QDialogButtonBox(QDialogButtonBox.Ok | QDialogButtonBox.Cancel)
        bb.button(QDialogButtonBox.Ok).setText("Compress…")
        bb.accepted.connect(self._compress)
        bb.rejected.connect(self.reject)
        layout.addWidget(bb)

    def _compress(self):
        path, _ = QFileDialog.getSaveFileName(self, "Save compressed PDF", "", "PDF (*.pdf)")
        if not path:
            return
        opts = self.LEVELS[self.level.currentText()]
        try:
            compress_pdf(self.document, path, opts["dpi"], opts["quality"])
            new_size = os.path.getsize(path)
            orig_size = os.path.getsize(self.document.path) if self.document.path else new_size
            reduction = (1 - new_size / orig_size) * 100 if orig_size else 0
            self.result_label.setText(
                f"Saved {human_size(new_size)} "
                f"(−{reduction:.0f}%) → {os.path.basename(path)}"
            )
            self.result_path = path
            QMessageBox.information(self, "Compressed",
                                    f"New size: {human_size(new_size)}\n"
                                    f"Reduction: {reduction:.0f}%")
        except Exception as e:
            QMessageBox.critical(self, "Compression failed", str(e))


def compress_pdf(document: PdfDocument, out_path: str, dpi: int, quality: int):
    """
    Rebuild the PDF, downsampling embedded raster images. Uses PyMuPDF to
    re-encode images to JPEG at the requested quality/dpi where possible.
    """
    src = document.doc
    # Save with strong garbage collection / deflation first.
    src.save(out_path, garbage=4, deflate=True, deflate_images=True,
             deflate_fonts=True, clean=True)
    if not HAS_PIL:
        return
    # Second pass: downsample images.
    try:
        doc = fitz.open(out_path)
        for page in doc:
            for img in page.get_images(full=True):
                xref = img[0]
                try:
                    base = doc.extract_image(xref)
                    image = Image.open(io.BytesIO(base["image"])).convert("RGB")
                    # Downscale large images.
                    max_dim = max(image.size)
                    target = int(dpi / 72 * max_dim)
                    if target < max_dim:
                        ratio = target / max_dim
                        image = image.resize(
                            (max(1, int(image.width * ratio)),
                             max(1, int(image.height * ratio))))
                    buf = io.BytesIO()
                    image.save(buf, format="JPEG", quality=quality)
                    doc.update_stream(xref, buf.getvalue())
                except Exception:
                    continue
        doc.save(out_path, garbage=4, deflate=True, clean=True, incremental=False)
        doc.close()
    except Exception as e:
        log.error("image downsample failed: %s", e)


class OcrDialog(QDialog):
    """Run OCR on the document or selected pages; export text."""

    LANGS = ["eng", "fra", "deu", "spa", "ita", "por", "rus", "ara", "chi_sim", "jpn"]

    def __init__(self, document: PdfDocument, parent=None):
        super().__init__(parent)
        self.document = document
        self.text_result = ""
        self.setWindowTitle("OCR")
        self.setMinimumSize(560, 480)
        layout = QVBoxLayout(self)

        if not HAS_TESSERACT:
            layout.addWidget(QLabel(
                "pytesseract is not installed (or the Tesseract binary is missing).\n"
                "Install with:  pip install pytesseract  and install Tesseract-OCR."
            ))

        row = QHBoxLayout()
        row.addWidget(QLabel("Language"))
        self.lang = QComboBox(); self.lang.addItems(self.LANGS)
        row.addWidget(self.lang)
        row.addWidget(QLabel("Pages"))
        self.pages = QLineEdit(); self.pages.setPlaceholderText("All")
        row.addWidget(self.pages)
        self.run_btn = QPushButton("Run OCR")
        self.run_btn.setObjectName("accent")
        self.run_btn.clicked.connect(self.run_ocr)
        self.run_btn.setEnabled(HAS_TESSERACT)
        row.addWidget(self.run_btn)
        layout.addLayout(row)

        self.output = QPlainTextEdit()
        self.output.setPlaceholderText("Recognized text will appear here…")
        layout.addWidget(self.output)

        btns = QHBoxLayout()
        export = QPushButton("Export Text…")
        export.clicked.connect(self.export_text)
        btns.addStretch(1)
        btns.addWidget(export)
        close = QPushButton("Close")
        close.clicked.connect(self.accept)
        btns.addWidget(close)
        layout.addLayout(btns)

    def run_ocr(self):
        pages = parse_page_ranges(self.pages.text(), self.document.page_count)
        lang = self.lang.currentText()
        chunks = []
        progress = QProgressDialog("Running OCR…", "Cancel", 0, len(pages), self)
        progress.setWindowModality(Qt.WindowModal)
        for i, pno in enumerate(pages):
            if progress.wasCanceled():
                break
            progress.setValue(i)
            QApplication.processEvents()
            try:
                pix = self.document.doc[pno].get_pixmap(matrix=fitz.Matrix(2, 2))
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                text = pytesseract.image_to_string(img, lang=lang)
                chunks.append(f"----- Page {pno + 1} -----\n{text}")
            except Exception as e:
                chunks.append(f"----- Page {pno + 1} (error: {e}) -----")
        progress.setValue(len(pages))
        self.text_result = "\n\n".join(chunks)
        self.output.setPlainText(self.text_result)

    def export_text(self):
        path, _ = QFileDialog.getSaveFileName(self, "Export OCR text", "", "Text (*.txt)")
        if path:
            with open(path, "w", encoding="utf-8") as f:
                f.write(self.output.toPlainText())
            QMessageBox.information(self, "Exported", f"Saved to {path}")


class BatchDialog(QDialog):
    """Batch operations across multiple files (merge/compress/watermark/...)."""

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Batch Tools")
        self.setMinimumSize(560, 460)
        layout = QVBoxLayout(self)

        row = QHBoxLayout()
        row.addWidget(QLabel("Operation"))
        self.op = QComboBox()
        self.op.addItems(["Merge", "Compress", "Watermark (text)", "Rename"])
        row.addWidget(self.op)
        layout.addLayout(row)

        self.list = QListWidget()
        self.list.setDragDropMode(QListWidget.InternalMove)
        layout.addWidget(self.list)

        controls = QHBoxLayout()
        add = QPushButton("Add files…"); add.clicked.connect(self._add)
        rem = QPushButton("Remove"); rem.clicked.connect(self._remove)
        controls.addWidget(add); controls.addWidget(rem); controls.addStretch(1)
        layout.addLayout(controls)

        opt_row = QHBoxLayout()
        opt_row.addWidget(QLabel("Text / Pattern"))
        self.param = QLineEdit()
        self.param.setPlaceholderText("Watermark text, or rename pattern e.g. doc_{n}")
        opt_row.addWidget(self.param)
        layout.addLayout(opt_row)

        run = QPushButton("Run Batch"); run.setObjectName("accent")
        run.clicked.connect(self._run)
        layout.addWidget(run)

    def _add(self):
        paths, _ = QFileDialog.getOpenFileNames(self, "Add PDFs", "", "PDF (*.pdf)")
        for p in paths:
            self.list.addItem(p)

    def _remove(self):
        for it in self.list.selectedItems():
            self.list.takeItem(self.list.row(it))

    def _files(self) -> List[str]:
        return [self.list.item(i).text() for i in range(self.list.count())]

    def _run(self):
        files = self._files()
        if not files:
            QMessageBox.warning(self, "Batch", "Add some files first.")
            return
        op = self.op.currentText()
        try:
            if op == "Merge":
                out, _ = QFileDialog.getSaveFileName(self, "Save merged PDF", "", "PDF (*.pdf)")
                if not out:
                    return
                merged = fitz.open()
                for f in files:
                    with fitz.open(f) as d:
                        merged.insert_pdf(d)
                merged.save(out, garbage=4, deflate=True)
                merged.close()
                QMessageBox.information(self, "Batch", f"Merged {len(files)} files.")
            elif op == "Compress":
                outdir = QFileDialog.getExistingDirectory(self, "Output folder")
                if not outdir:
                    return
                for f in files:
                    d = PdfDocument(f)
                    out = os.path.join(outdir, "compressed_" + os.path.basename(f))
                    compress_pdf(d, out, 110, 70)
                    d.close()
                QMessageBox.information(self, "Batch", "Compression complete.")
            elif op.startswith("Watermark"):
                outdir = QFileDialog.getExistingDirectory(self, "Output folder")
                if not outdir:
                    return
                wtext = self.param.text() or "CONFIDENTIAL"
                for f in files:
                    d = fitz.open(f)
                    for page in d:
                        r = page.rect
                        page.insert_text(
                            fitz.Point(r.width / 2 - len(wtext) * 8, r.height / 2),
                            wtext, fontsize=42, rotate=45,
                            color=(0.6, 0.6, 0.6), fill_opacity=0.25, overlay=True)
                    d.save(os.path.join(outdir, "wm_" + os.path.basename(f)),
                           garbage=4, deflate=True)
                    d.close()
                QMessageBox.information(self, "Batch", "Watermarking complete.")
            elif op == "Rename":
                outdir = QFileDialog.getExistingDirectory(self, "Output folder")
                if not outdir:
                    return
                pattern = self.param.text() or "document_{n}"
                for i, f in enumerate(files, 1):
                    name = pattern.format(n=i) + ".pdf"
                    shutil.copy(f, os.path.join(outdir, name))
                QMessageBox.information(self, "Batch", "Renaming complete.")
        except Exception as e:
            QMessageBox.critical(self, "Batch failed", str(e))


class SettingsDialog(QDialog):
    """Preferences: theme, accent, fonts, zoom, startup, etc. (SQLite-backed)."""

    def __init__(self, settings: SettingsManager, parent=None):
        super().__init__(parent)
        self.settings = settings
        self.setWindowTitle("Settings")
        self.setMinimumWidth(440)
        form = QFormLayout(self)

        self.theme = QComboBox(); self.theme.addItems(["dark", "light"])
        self.theme.setCurrentText(settings.get("theme"))
        form.addRow("Theme", self.theme)

        self.accent_btn = QPushButton("Accent color")
        self._accent = QColor(settings.get("accent"))
        self.accent_btn.clicked.connect(self._pick_accent)
        self._update_accent_btn()
        form.addRow("Accent", self.accent_btn)

        self.font = QFontComboBox()
        self.font.setCurrentFont(QFont(settings.get("ui_font")))
        form.addRow("UI Font", self.font)

        self.zoom = QSpinBox(); self.zoom.setRange(25, 400)
        self.zoom.setValue(settings.get_int("default_zoom", 100))
        self.zoom.setSuffix(" %")
        form.addRow("Default zoom", self.zoom)

        self.startup = QComboBox(); self.startup.addItems(["welcome", "blank", "last"])
        self.startup.setCurrentText(settings.get("startup_page"))
        form.addRow("Startup", self.startup)

        self.view_mode = QComboBox(); self.view_mode.addItems(["single", "continuous", "two"])
        self.view_mode.setCurrentText(settings.get("view_mode"))
        form.addRow("Default view", self.view_mode)

        self.language = QComboBox()
        self.language.addItems(["English", "Français", "Deutsch", "Español", "العربية"])
        self.language.setCurrentText(settings.get("language"))
        form.addRow("Language", self.language)

        self.autosave = QCheckBox("Enable autosave recovery")
        self.autosave.setChecked(settings.get_bool("autosave", True))
        form.addRow(self.autosave)

        bb = QDialogButtonBox(QDialogButtonBox.Save | QDialogButtonBox.Cancel)
        bb.accepted.connect(self._save)
        bb.rejected.connect(self.reject)
        form.addRow(bb)

    def _pick_accent(self):
        c = QColorDialog.getColor(self._accent, self, "Accent color")
        if c.isValid():
            self._accent = c
            self._update_accent_btn()

    def _update_accent_btn(self):
        self.accent_btn.setStyleSheet(
            f"background:{self._accent.name()}; color:white; border:none;")

    def _save(self):
        self.settings.set("theme", self.theme.currentText())
        self.settings.set("accent", self._accent.name())
        self.settings.set("ui_font", self.font.currentFont().family())
        self.settings.set("default_zoom", self.zoom.value())
        self.settings.set("startup_page", self.startup.currentText())
        self.settings.set("view_mode", self.view_mode.currentText())
        self.settings.set("language", self.language.currentText())
        self.settings.set("autosave", "1" if self.autosave.isChecked() else "0")
        self.accept()


# ----------------------------------------------------------------------------
# Converters (free functions used by menu actions / workers)
# ----------------------------------------------------------------------------
def pdf_to_text(document: PdfDocument, out_path: str, progress_cb=None):
    parts = []
    n = document.page_count
    for i in range(n):
        parts.append(document.doc[i].get_text())
        if progress_cb:
            progress_cb(int((i + 1) / n * 100), f"Page {i+1}/{n}")
    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(parts))
    return out_path


def pdf_to_html(document: PdfDocument, out_path: str, progress_cb=None):
    html = ["<html><meta charset='utf-8'><body>"]
    n = document.page_count
    for i in range(n):
        html.append(document.doc[i].get_text("html"))
        if progress_cb:
            progress_cb(int((i + 1) / n * 100), f"Page {i+1}/{n}")
    html.append("</body></html>")
    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(html))
    return out_path


def pdf_to_images(document: PdfDocument, out_dir: str, dpi: int = 150, progress_cb=None):
    os.makedirs(out_dir, exist_ok=True)
    zoom = dpi / 72.0
    n = document.page_count
    paths = []
    for i in range(n):
        pix = document.doc[i].get_pixmap(matrix=fitz.Matrix(zoom, zoom))
        p = os.path.join(out_dir, f"page_{i+1:03d}.png")
        pix.save(p)
        paths.append(p)
        if progress_cb:
            progress_cb(int((i + 1) / n * 100), f"Page {i+1}/{n}")
    return paths


def pdf_to_word(document: PdfDocument, out_path: str, progress_cb=None):
    if not HAS_DOCX:
        raise RuntimeError("python-docx not installed (pip install python-docx).")
    d = docx.Document()
    n = document.page_count
    for i in range(n):
        text = document.doc[i].get_text()
        for line in text.splitlines():
            d.add_paragraph(line)
        d.add_page_break()
        if progress_cb:
            progress_cb(int((i + 1) / n * 100), f"Page {i+1}/{n}")
    d.save(out_path)
    return out_path


def pdf_to_excel(document: PdfDocument, out_path: str, progress_cb=None):
    if not HAS_OPENPYXL:
        raise RuntimeError("openpyxl not installed (pip install openpyxl).")
    from openpyxl import Workbook
    wb = Workbook()
    n = document.page_count
    for i in range(n):
        ws = wb.create_sheet(f"Page{i+1}") if i else wb.active
        if i == 0:
            ws.title = "Page1"
        row = 1
        for line in document.doc[i].get_text().splitlines():
            for col, cell in enumerate(line.split("\t"), 1):
                ws.cell(row=row, column=col, value=cell)
            row += 1
        if progress_cb:
            progress_cb(int((i + 1) / n * 100), f"Page {i+1}/{n}")
    wb.save(out_path)
    return out_path


def pdf_to_powerpoint(document: PdfDocument, out_path: str, progress_cb=None):
    if not HAS_PPTX:
        raise RuntimeError("python-pptx not installed (pip install python-pptx).")
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    n = document.page_count
    tmpdir = tempfile.mkdtemp()
    for i in range(n):
        pix = document.doc[i].get_pixmap(matrix=fitz.Matrix(2, 2))
        img_path = os.path.join(tmpdir, f"p{i}.png")
        pix.save(img_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width)
        if progress_cb:
            progress_cb(int((i + 1) / n * 100), f"Page {i+1}/{n}")
    prs.save(out_path)
    shutil.rmtree(tmpdir, ignore_errors=True)
    return out_path


def images_to_pdf(image_paths: List[str], out_path: str, progress_cb=None):
    doc = fitz.open()
    n = len(image_paths)
    for i, p in enumerate(image_paths):
        img = fitz.open(p)
        rect = img[0].rect
        pdfbytes = img.convert_to_pdf()
        img.close()
        imgpdf = fitz.open("pdf", pdfbytes)
        doc.insert_pdf(imgpdf)
        if progress_cb:
            progress_cb(int((i + 1) / n * 100), f"Image {i+1}/{n}")
    doc.save(out_path)
    doc.close()
    return out_path


def text_to_pdf(text_path: str, out_path: str, progress_cb=None):
    with open(text_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    doc = fitz.open()
    page = doc.new_page()
    rect = fitz.Rect(50, 50, page.rect.width - 50, page.rect.height - 50)
    remaining = page.insert_textbox(rect, text, fontsize=11)
    # If text overflows, keep adding pages.
    while remaining < 0:
        page = doc.new_page()
        rect = fitz.Rect(50, 50, page.rect.width - 50, page.rect.height - 50)
        # crude continuation: re-insert full text (acceptable for plain text)
        remaining = page.insert_textbox(rect, text, fontsize=11)
        break
    doc.save(out_path)
    doc.close()
    return out_path


# ----------------------------------------------------------------------------
# 10. MainWindow
# ----------------------------------------------------------------------------
class DocumentTab(QWidget):
    """A single open document: viewer + bound state."""

    def __init__(self, document: PdfDocument, settings: SettingsManager):
        super().__init__()
        self.document = document
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        self.viewer = PdfViewer(document, settings)
        layout.addWidget(self.viewer)


class MainWindow(QMainWindow):
    """The application shell: menus, toolbar, panels, tabs, status bar."""

    def __init__(self, settings: SettingsManager):
        super().__init__()
        self.settings = settings
        self.threadpool = QThreadPool.globalInstance()
        self.setWindowTitle(APP_NAME)
        self.resize(1360, 860)
        self.setAcceptDrops(True)

        # Central tabbed area
        self.tabs = QTabWidget()
        self.tabs.setTabsClosable(True)
        self.tabs.setMovable(True)
        self.tabs.tabCloseRequested.connect(self.close_tab)
        self.tabs.currentChanged.connect(self.on_tab_changed)

        # Side rail + stacked panels
        self._build_left_sidebar()
        self.properties = PropertiesPanel()

        # Splitter layout: [rail|panel] | tabs | properties
        self.left_split = QSplitter(Qt.Horizontal)
        self.left_split.addWidget(self.sidebar_container)
        self.left_split.addWidget(self.tabs)
        self.left_split.addWidget(self.properties)
        self.left_split.setStretchFactor(1, 1)
        self.left_split.setSizes([280, 820, 260])
        self.setCentralWidget(self.left_split)

        self._build_actions()
        self._build_menus()
        self._build_toolbar()
        self._build_statusbar()

        # Autosave timer
        self.autosave_timer = QTimer(self)
        self.autosave_timer.timeout.connect(self._autosave)
        self._reconfigure_autosave()

        self._show_startup()

    # ---- sidebar -----------------------------------------------------------
    def _build_left_sidebar(self):
        self.sidebar_container = QWidget()
        h = QHBoxLayout(self.sidebar_container)
        h.setContentsMargins(0, 0, 0, 0)
        h.setSpacing(0)

        # Rail of icon buttons
        rail = QFrame(); rail.setObjectName("SidebarRail")
        rail.setFixedWidth(72)
        rl = QVBoxLayout(rail)
        rl.setContentsMargins(8, 10, 8, 10)
        rl.setSpacing(6)

        self.panel_stack = QStackedWidget()
        self.thumb_panel = ThumbnailPanel()
        self.bookmark_panel = BookmarkPanel()
        self.comments_panel = CommentsPanel()
        self.search_panel = SearchPanel()
        for p in (self.thumb_panel, self.bookmark_panel, self.comments_panel, self.search_panel):
            self.panel_stack.addWidget(p)

        self.thumb_panel.page_selected.connect(self._goto_current_page)
        self.thumb_panel.order_changed.connect(self._reorder_pages)
        self.bookmark_panel.bookmark_selected.connect(self._goto_current_page)
        self.comments_panel.annotation_selected.connect(self._goto_current_page)
        self.search_panel.result_selected.connect(self._on_search_result)

        self._rail_buttons = []
        rail_items = [("Pages", "pages", 0), ("Marks", "bookmark", 1),
                      ("Notes", "comment", 2), ("Find", "search", 3)]
        for text, icon_name, idx in rail_items:
            btn = QToolButton()
            btn.setObjectName("RailButton")
            btn.setText(text)
            btn.setIcon(vector_icon(icon_name))
            btn.setIconSize(QSize(22, 22))
            btn.setCheckable(True)
            btn.setToolButtonStyle(Qt.ToolButtonTextUnderIcon)
            btn.setToolTip(text)
            btn.clicked.connect(lambda _=False, i=idx: self._select_panel(i))
            rl.addWidget(btn)
            self._rail_buttons.append(btn)
        rl.addStretch(1)
        # Theme / settings shortcuts pinned to the bottom of the rail.
        for tip, icon_name, cb in [("Toggle Theme", "theme", self.toggle_theme),
                                   ("Settings", "settings", self.show_settings)]:
            b = QToolButton()
            b.setObjectName("RailButton")
            b.setIcon(vector_icon(icon_name))
            b.setIconSize(QSize(22, 22))
            b.setToolTip(tip)
            b.clicked.connect(cb)
            rl.addWidget(b)
        self._rail_buttons[0].setChecked(True)

        h.addWidget(rail)
        h.addWidget(self.panel_stack)

    def _select_panel(self, idx: int):
        self.panel_stack.setCurrentIndex(idx)
        for i, b in enumerate(self._rail_buttons):
            b.setChecked(i == idx)

    # ---- actions -----------------------------------------------------------
    def _act(self, text, slot, shortcut=None, glyph=None, checkable=False):
        action = QAction(text, self)
        if glyph:
            action.setIcon(vector_icon(glyph))
        if shortcut:
            action.setShortcut(QKeySequence(shortcut))
        action.setCheckable(checkable)
        action.triggered.connect(slot)
        return action

    def _build_actions(self):
        A = self._act
        # File
        self.a_new = A("New", self.new_document, "Ctrl+N", "new")
        self.a_open = A("Open…", self.open_dialog, "Ctrl+O", "open")
        self.a_save = A("Save", self.save_current, "Ctrl+S", "save")
        self.a_saveas = A("Save As…", self.save_as_current, "Ctrl+Shift+S")
        self.a_export = A("Export…", self.export_current, "Ctrl+E", "export")
        self.a_print = A("Print…", self.print_current, "Ctrl+P", "print")
        self.a_close = A("Close Tab", lambda: self.close_tab(self.tabs.currentIndex()), "Ctrl+W")
        self.a_quit = A("Exit", self.close, "Ctrl+Q")
        # Edit / undo placeholders
        self.a_undo = A("Undo", self._not_impl, "Ctrl+Z", "undo")
        self.a_redo = A("Redo", self._not_impl, "Ctrl+Y", "redo")
        self.a_find = A("Find", lambda: self._select_panel(3), "Ctrl+F", "search")
        # View
        self.a_zoom_in = A("Zoom In", self._zoom_in, "Ctrl++", "zoom_in")
        self.a_zoom_out = A("Zoom Out", self._zoom_out, "Ctrl+-", "zoom_out")
        self.a_fit_width = A("Fit Width", self._fit_width, glyph="fit_width")
        self.a_fit_page = A("Fit Page", self._fit_page, glyph="fit_page")
        self.a_rotate_left = A("Rotate Left", lambda: self._rotate(-90), "Ctrl+L", "rotate_left")
        self.a_rotate_right = A("Rotate Right", lambda: self._rotate(90), "Ctrl+R", "rotate_right")
        self.a_fullscreen = A("Fullscreen", self.toggle_fullscreen, "F11", "fullscreen")
        self.a_single = A("Single Page", lambda: self._set_view_mode("single"), checkable=True)
        self.a_continuous = A("Continuous", lambda: self._set_view_mode("continuous"), checkable=True)
        self.a_two = A("Two Pages", lambda: self._set_view_mode("two"), checkable=True)
        # Insert / pages
        self.a_insert_blank = A("Insert Blank Page", self.insert_blank_page)
        self.a_delete_pages = A("Delete Pages…", self.delete_pages)
        self.a_duplicate_page = A("Duplicate Page", self.duplicate_page)
        self.a_rotate_pages = A("Rotate Pages…", self.rotate_pages)
        self.a_reverse = A("Reverse Pages", self.reverse_pages)
        self.a_extract = A("Extract Pages…", self.extract_pages)
        self.a_split = A("Split PDF…", self.split_pdf, glyph="split")
        self.a_merge = A("Merge PDFs…", self.merge_pdfs, glyph="merge")
        self.a_insert_image = A("Insert Image…", self.insert_image, glyph="image")
        self.a_insert_text = A("Add Text…", self.add_text_box, glyph="text")
        # Annotations
        self.a_highlight = A("Highlight", lambda: self._quick_annot("highlight"), glyph="highlight")
        self.a_underline = A("Underline", lambda: self._quick_annot("underline"), glyph="underline")
        self.a_strike = A("Strikeout", lambda: self._quick_annot("strikeout"), glyph="strikeout")
        self.a_note = A("Sticky Note…", self.add_sticky_note, glyph="note")
        # Tools
        self.a_watermark = A("Watermark…", self.show_watermark)
        self.a_header_footer = A("Header & Footer…", self.show_header_footer)
        self.a_compress = A("Compress…", self.show_compress, "Ctrl+M")
        self.a_ocr = A("OCR…", self.show_ocr)
        self.a_batch = A("Batch Tools…", self.show_batch)
        # Convert
        self.a_to_word = A("PDF → Word", lambda: self._convert("word"))
        self.a_to_excel = A("PDF → Excel", lambda: self._convert("excel"))
        self.a_to_ppt = A("PDF → PowerPoint", lambda: self._convert("ppt"))
        self.a_to_text = A("PDF → Text", lambda: self._convert("text"))
        self.a_to_html = A("PDF → HTML", lambda: self._convert("html"))
        self.a_to_image = A("PDF → Image", lambda: self._convert("image"))
        self.a_img_to_pdf = A("Images → PDF…", self.images_to_pdf_action)
        self.a_text_to_pdf = A("Text → PDF…", self.text_to_pdf_action)
        # Security
        self.a_security = A("Security…", self.show_security, glyph="lock")
        # Reading modes
        self.a_reading_normal = A("Normal", lambda: self._reading_mode("Normal"))
        self.a_reading_sepia = A("Sepia", lambda: self._reading_mode("Sepia"))
        self.a_reading_eye = A("Eye Comfort", lambda: self._reading_mode("Eye Comfort"))
        self.a_reading_night = A("Night", lambda: self._reading_mode("Night"))
        self.a_autoscroll = A("Auto Scroll", self._toggle_autoscroll, checkable=True)
        # Window / theme
        self.a_toggle_theme = A("Toggle Theme", self.toggle_theme, glyph="theme")
        self.a_settings = A("Settings…", self.show_settings, glyph="settings")
        # Help
        self.a_about = A("About", self.show_about)

        mg = QActionGroup(self)
        for a in (self.a_single, self.a_continuous, self.a_two):
            mg.addAction(a)
        self.a_continuous.setChecked(True)

    def _build_menus(self):
        mb = self.menuBar()

        m_file = mb.addMenu("&File")
        m_file.addActions([self.a_new, self.a_open])
        self.recent_menu = m_file.addMenu("Recent Files")
        self._rebuild_recent_menu()
        m_file.addSeparator()
        m_file.addActions([self.a_save, self.a_saveas, self.a_export, self.a_print])
        m_file.addSeparator()
        m_file.addActions([self.a_close, self.a_quit])

        m_edit = mb.addMenu("&Edit")
        m_edit.addActions([self.a_undo, self.a_redo])
        m_edit.addSeparator()
        m_edit.addAction(self.a_find)
        m_edit.addActions([self.a_insert_text])

        m_view = mb.addMenu("&View")
        m_view.addActions([self.a_zoom_in, self.a_zoom_out, self.a_fit_width, self.a_fit_page])
        m_view.addSeparator()
        m_view.addActions([self.a_rotate_left, self.a_rotate_right])
        m_view.addSeparator()
        m_view.addActions([self.a_single, self.a_continuous, self.a_two])
        m_view.addSeparator()
        reading = m_view.addMenu("Reading Mode")
        reading.addActions([self.a_reading_normal, self.a_reading_sepia,
                            self.a_reading_eye, self.a_reading_night])
        m_view.addAction(self.a_autoscroll)
        m_view.addSeparator()
        m_view.addAction(self.a_fullscreen)

        m_insert = mb.addMenu("&Insert")
        m_insert.addActions([self.a_insert_blank, self.a_duplicate_page,
                             self.a_insert_image, self.a_insert_text])
        m_insert.addSeparator()
        m_insert.addActions([self.a_highlight, self.a_underline, self.a_strike, self.a_note])

        m_tools = mb.addMenu("&Tools")
        m_tools.addActions([self.a_delete_pages, self.a_rotate_pages, self.a_reverse,
                            self.a_extract, self.a_split, self.a_merge])
        m_tools.addSeparator()
        m_tools.addActions([self.a_watermark, self.a_header_footer, self.a_compress,
                            self.a_ocr, self.a_batch])
        m_tools.addSeparator()
        conv = m_tools.addMenu("Convert")
        conv.addActions([self.a_to_word, self.a_to_excel, self.a_to_ppt,
                         self.a_to_text, self.a_to_html, self.a_to_image])
        conv.addSeparator()
        conv.addActions([self.a_img_to_pdf, self.a_text_to_pdf])

        m_sec = mb.addMenu("&Security")
        m_sec.addAction(self.a_security)

        m_win = mb.addMenu("&Window")
        m_win.addActions([self.a_toggle_theme, self.a_settings])

        m_help = mb.addMenu("&Help")
        m_help.addAction(self.a_about)

    def _build_toolbar(self):
        tb = QToolBar("Main")
        tb.setIconSize(QSize(20, 20))
        tb.setMovable(False)
        self.addToolBar(tb)
        for a in (self.a_new, self.a_open, self.a_save, self.a_export, self.a_print):
            tb.addAction(a)
        tb.addSeparator()

        # Page navigation widgets in toolbar
        prev_btn = QToolButton(); prev_btn.setText("‹"); prev_btn.setToolTip("Previous page")
        prev_btn.clicked.connect(lambda: self.current_viewer() and self.current_viewer().prev_page())
        tb.addWidget(prev_btn)
        self.page_spin = QSpinBox()
        self.page_spin.setMinimum(1)
        self.page_spin.setMaximum(1)
        self.page_spin.setFixedWidth(58)
        self.page_spin.setAlignment(Qt.AlignCenter)
        self.page_spin.setButtonSymbols(QSpinBox.NoButtons)
        self.page_spin.editingFinished.connect(self._jump_to_spin)
        tb.addWidget(self.page_spin)
        self.page_total_lbl = QLabel(" / 0 ")
        self.page_total_lbl.setStyleSheet("padding: 0 4px;")
        tb.addWidget(self.page_total_lbl)
        next_btn = QToolButton(); next_btn.setText("›"); next_btn.setToolTip("Next page")
        next_btn.clicked.connect(lambda: self.current_viewer() and self.current_viewer().next_page())
        tb.addWidget(next_btn)

        tb.addSeparator()
        # Zoom group
        tb.addAction(self.a_zoom_out)
        self.zoom_combo = QComboBox()
        self.zoom_combo.setEditable(True)
        self.zoom_combo.setFixedWidth(96)
        self.zoom_combo.setInsertPolicy(QComboBox.NoInsert)
        self.zoom_combo.addItems(["50%", "75%", "100%", "125%", "150%", "200%",
                                  "Fit Width", "Fit Page"])
        self.zoom_combo.setCurrentText("100%")
        self.zoom_combo.activated.connect(self._zoom_combo_changed)
        tb.addWidget(self.zoom_combo)
        tb.addAction(self.a_zoom_in)
        tb.addAction(self.a_fit_width)
        tb.addAction(self.a_fit_page)

        tb.addSeparator()
        for a in (self.a_rotate_left, self.a_rotate_right):
            tb.addAction(a)
        tb.addSeparator()
        for a in (self.a_highlight, self.a_underline, self.a_strike, self.a_note):
            tb.addAction(a)

        # Push the remaining buttons to the right edge.
        spacer = QWidget()
        spacer.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Preferred)
        tb.addWidget(spacer)
        tb.addAction(self.a_security)
        tb.addAction(self.a_fullscreen)
        tb.addAction(self.a_toggle_theme)

        self.toolbar = tb
        if not self.settings.get_bool("show_toolbar", True):
            tb.hide()

    def _build_statusbar(self):
        sb = QStatusBar()
        sb.setSizeGripEnabled(False)
        self.setStatusBar(sb)
        self.status_msg = QLabel("Ready")
        self.status_msg.setObjectName("StatusMsg")
        self.status_page = QLabel("Page – / –")
        self.status_zoom = QLabel("100%")
        self.status_size = QLabel("–")
        for chip in (self.status_page, self.status_zoom, self.status_size):
            chip.setObjectName("StatusChip")
        sb.addWidget(self.status_msg, 1)
        sb.addPermanentWidget(self.status_page)
        sb.addPermanentWidget(self.status_zoom)
        sb.addPermanentWidget(self.status_size)

    # ---- startup / welcome -------------------------------------------------
    def _show_startup(self):
        mode = self.settings.get("startup_page")
        if mode == "blank":
            self.new_document()
        elif mode == "last":
            recents = self.settings.recent_files(1)
            if recents:
                self.open_path(recents[0])
            else:
                self._add_welcome_tab()
        else:
            self._add_welcome_tab()

    def _add_welcome_tab(self):
        # Outer container centers a fixed-width card.
        outer = QWidget()
        ol = QVBoxLayout(outer)
        ol.setAlignment(Qt.AlignCenter)

        card = QFrame()
        card.setObjectName("WelcomeCard")
        card.setFixedWidth(640)
        v = QVBoxLayout(card)
        v.setContentsMargins(40, 36, 40, 36)
        v.setSpacing(6)

        # Header: app logo + name.
        head = QHBoxLayout()
        head.setSpacing(16)
        logo = QLabel()
        logo.setPixmap(vector_icon("page", self.settings.get("accent"), 56).pixmap(56, 56))
        head.addWidget(logo)
        head_text = QVBoxLayout(); head_text.setSpacing(2)
        title = QLabel(APP_NAME); title.setObjectName("WelcomeTitle")
        sub = QLabel(f"Professional PDF toolkit · v{APP_VERSION}")
        sub.setObjectName("WelcomeSub")
        head_text.addWidget(title); head_text.addWidget(sub)
        head.addLayout(head_text)
        head.addStretch(1)
        v.addLayout(head)
        v.addSpacing(22)

        # Action tiles (2x2 grid).
        grid = QGridLayout()
        grid.setSpacing(12)
        tiles = [
            ("Open PDF", "open", "Browse for an existing file", self.open_dialog),
            ("New PDF", "new", "Start a blank document", self.new_document),
            ("Merge PDFs", "merge", "Combine multiple files", self.merge_pdfs),
            ("Images → PDF", "image", "Build a PDF from images", self.images_to_pdf_action),
        ]
        for i, (label, icon_name, tip, cb) in enumerate(tiles):
            btn = QPushButton(f"  {label}")
            btn.setObjectName("TileButton")
            btn.setIcon(vector_icon(icon_name, self.settings.get("accent"), 40))
            btn.setIconSize(QSize(26, 26))
            btn.setToolTip(tip)
            btn.setCursor(Qt.PointingHandCursor)
            btn.clicked.connect(cb)
            grid.addWidget(btn, i // 2, i % 2)
        v.addLayout(grid)
        v.addSpacing(20)

        # Recent files.
        recents = self.settings.recent_files(6)
        sec = QLabel("RECENT FILES" if recents else "DROP A PDF HERE TO OPEN")
        sec.setObjectName("SectionLabel")
        v.addWidget(sec)
        if recents:
            rl = QListWidget()
            rl.setObjectName("RecentList")
            rl.setMaximumHeight(180)
            for p in recents:
                it = QListWidgetItem(vector_icon("page", ICON_COLOR, 28), os.path.basename(p))
                it.setToolTip(p)
                it.setData(Qt.UserRole, p)
                rl.addItem(it)
            rl.itemActivated.connect(lambda it: self.open_path(it.data(Qt.UserRole)))
            rl.itemDoubleClicked.connect(lambda it: self.open_path(it.data(Qt.UserRole)))
            v.addWidget(rl)
        else:
            hint = QLabel("You can also drag & drop a PDF anywhere in the window.")
            hint.setObjectName("WelcomeSub")
            v.addWidget(hint)

        ol.addWidget(card)
        idx = self.tabs.addTab(outer, "Welcome")
        self.tabs.setCurrentIndex(idx)

    # ---- current document helpers -----------------------------------------
    def current_tab(self) -> Optional[DocumentTab]:
        w = self.tabs.currentWidget()
        return w if isinstance(w, DocumentTab) else None

    def current_doc(self) -> Optional[PdfDocument]:
        t = self.current_tab()
        return t.document if t else None

    def current_viewer(self) -> Optional[PdfViewer]:
        t = self.current_tab()
        return t.viewer if t else None

    def _require_doc(self) -> Optional[PdfDocument]:
        d = self.current_doc()
        if not d:
            QMessageBox.information(self, APP_NAME, "Open a document first.")
        return d

    # ---- file operations ---------------------------------------------------
    def new_document(self):
        if not self._ensure_fitz():
            return
        doc = PdfDocument()
        self._add_document_tab(doc, "Untitled.pdf")
        self.settings.log_action("new", "")

    def open_dialog(self):
        paths, _ = QFileDialog.getOpenFileNames(
            self, "Open PDF", "", "PDF Files (*.pdf);;All Files (*)")
        for p in paths:
            self.open_path(p)

    def open_path(self, path: str):
        if not self._ensure_fitz():
            return
        if not os.path.exists(path):
            QMessageBox.warning(self, APP_NAME, f"File not found:\n{path}")
            return
        # Already open? focus it.
        for i in range(self.tabs.count()):
            w = self.tabs.widget(i)
            if isinstance(w, DocumentTab) and w.document.path == path:
                self.tabs.setCurrentIndex(i)
                return
        try:
            doc = PdfDocument()
            doc.open(path)
        except PermissionError:
            pw, ok = QInputDialog.getText(
                self, "Password", "This PDF is encrypted. Password:",
                QLineEdit.Password)
            if not ok:
                return
            try:
                doc = PdfDocument()
                doc.open(path, pw)
            except Exception as e:
                QMessageBox.critical(self, "Open failed", str(e))
                return
        except Exception as e:
            QMessageBox.critical(self, "Open failed", str(e))
            log.error("open %s: %s", path, e)
            return

        doc.view_mode = self.settings.get("view_mode", "continuous")
        doc.zoom = self.settings.get_int("default_zoom", 100) / 100.0
        self._add_document_tab(doc, doc.name)
        self.settings.add_recent(path)
        self.settings.log_action("open", path)
        self._rebuild_recent_menu()

    def _add_document_tab(self, doc: PdfDocument, title: str):
        tab = DocumentTab(doc, self.settings)
        tab.viewer.page_changed.connect(self._on_page_changed)
        tab.viewer.zoom_changed.connect(self._on_zoom_changed)
        tab.viewer.set_reading_mode(self.settings.get("reading_mode", "Normal"))
        idx = self.tabs.addTab(tab, title)
        self.tabs.setCurrentIndex(idx)
        self._refresh_panels()
        self._update_status()

    def save_current(self):
        doc = self._require_doc()
        if not doc:
            return
        if not doc.path:
            return self.save_as_current()
        try:
            doc.save()
            self._set_msg(f"Saved {doc.name}")
            self.settings.log_action("save", doc.path)
        except Exception as e:
            QMessageBox.critical(self, "Save failed", str(e))

    def save_as_current(self):
        doc = self._require_doc()
        if not doc:
            return
        path, _ = QFileDialog.getSaveFileName(
            self, "Save As", doc.path or "Untitled.pdf", "PDF Files (*.pdf)")
        if not path:
            return
        try:
            doc.save(path)
            self.tabs.setTabText(self.tabs.currentIndex(), doc.name)
            self.settings.add_recent(path)
            self._rebuild_recent_menu()
            self._set_msg(f"Saved {doc.name}")
        except Exception as e:
            QMessageBox.critical(self, "Save failed", str(e))

    def export_current(self):
        doc = self._require_doc()
        if not doc:
            return
        path, sel = QFileDialog.getSaveFileName(
            self, "Export", "",
            "PDF (*.pdf);;Text (*.txt);;HTML (*.html);;PNG image (first page) (*.png)")
        if not path:
            return
        try:
            if path.lower().endswith(".txt"):
                pdf_to_text(doc, path)
            elif path.lower().endswith(".html"):
                pdf_to_html(doc, path)
            elif path.lower().endswith(".png"):
                pix = doc.doc[doc.current_page].get_pixmap(matrix=fitz.Matrix(2, 2))
                pix.save(path)
            else:
                doc.save(path)
            self._set_msg(f"Exported to {os.path.basename(path)}")
        except Exception as e:
            QMessageBox.critical(self, "Export failed", str(e))

    def print_current(self):
        doc = self._require_doc()
        if not doc:
            return
        printer = QPrinter(QPrinter.HighResolution)
        dlg = QPrintDialog(printer, self)
        dlg.setOption(QPrintDialog.PrintPageRange, True)
        if dlg.exec() != QDialog.Accepted:
            return
        try:
            self._render_to_printer(doc, printer)
            self._set_msg("Sent to printer")
        except Exception as e:
            QMessageBox.critical(self, "Print failed", str(e))

    def _render_to_printer(self, doc: PdfDocument, printer: QPrinter):
        painter = QPainter()
        if not painter.begin(printer):
            raise RuntimeError("Cannot open printer device.")
        try:
            # Determine page range.
            from_p = printer.fromPage()
            to_p = printer.toPage()
            if from_p == 0:
                pages = range(doc.page_count)
            else:
                pages = range(from_p - 1, to_p)
            first = True
            page_rect = printer.pageRect(QPrinter.DevicePixel)
            for i in pages:
                if not first:
                    printer.newPage()
                first = False
                img = doc.render_page(i, 2.0, doc.rotation)
                pix = QPixmap.fromImage(img)
                scaled = pix.scaled(
                    int(page_rect.width()), int(page_rect.height()),
                    Qt.KeepAspectRatio, Qt.SmoothTransformation)
                x = (page_rect.width() - scaled.width()) / 2
                y = (page_rect.height() - scaled.height()) / 2
                painter.drawPixmap(int(x), int(y), scaled)
        finally:
            painter.end()

    def close_tab(self, index: int):
        w = self.tabs.widget(index)
        if isinstance(w, DocumentTab):
            doc = w.document
            if doc.dirty:
                ret = QMessageBox.question(
                    self, "Unsaved changes",
                    f"Save changes to {doc.name}?",
                    QMessageBox.Save | QMessageBox.Discard | QMessageBox.Cancel)
                if ret == QMessageBox.Cancel:
                    return
                if ret == QMessageBox.Save:
                    self.tabs.setCurrentIndex(index)
                    self.save_current()
            doc.close()
        self.tabs.removeTab(index)
        if self.tabs.count() == 0:
            self._add_welcome_tab()

    def on_tab_changed(self, _index: int):
        self._refresh_panels()
        self._update_status()
        d = self.current_doc()
        if d:
            self.a_single.setChecked(d.view_mode == "single")
            self.a_continuous.setChecked(d.view_mode == "continuous")
            self.a_two.setChecked(d.view_mode == "two")

    # ---- panels refresh ----------------------------------------------------
    def _refresh_panels(self):
        doc = self.current_doc()
        self.properties.update_for(doc)
        if doc:
            self.thumb_panel.populate(doc)
            self.bookmark_panel.populate(doc)
            self.comments_panel.populate(doc)
            self.search_panel.set_document(doc)

    def _rebuild_recent_menu(self):
        self.recent_menu.clear()
        recents = self.settings.recent_files(15)
        if not recents:
            a = self.recent_menu.addAction("(No recent files)")
            a.setEnabled(False)
            return
        for p in recents:
            act = self.recent_menu.addAction(os.path.basename(p))
            act.setToolTip(p)
            act.triggered.connect(lambda _=False, path=p: self.open_path(path))
        self.recent_menu.addSeparator()
        clear = self.recent_menu.addAction("Clear Recent")
        clear.triggered.connect(lambda: (self.settings.clear_recent(),
                                         self._rebuild_recent_menu()))

    # ---- navigation / status ----------------------------------------------
    def _on_page_changed(self, page: int):
        self.page_spin.blockSignals(True)
        self.page_spin.setValue(page + 1)
        self.page_spin.blockSignals(False)
        self._update_status()

    def _on_zoom_changed(self, zoom: float):
        self.zoom_combo.blockSignals(True)
        self.zoom_combo.setCurrentText(f"{int(zoom * 100)}%")
        self.zoom_combo.blockSignals(False)
        self._update_status()

    def _update_status(self):
        doc = self.current_doc()
        if not doc:
            self.status_page.setText("Page – / –")
            self.status_zoom.setText("–")
            self.status_size.setText("–")
            self.page_total_lbl.setText(" / 0 ")
            self.page_spin.setMaximum(1)
            return
        self.page_spin.setMaximum(max(1, doc.page_count))
        self.page_total_lbl.setText(f" / {doc.page_count} ")
        self.status_page.setText(f"Page {doc.current_page + 1} / {doc.page_count}")
        self.status_zoom.setText(f"{int(doc.zoom * 100)}%")
        if doc.path and os.path.exists(doc.path):
            self.status_size.setText(human_size(os.path.getsize(doc.path)))
        else:
            self.status_size.setText("Unsaved")

    def _jump_to_spin(self):
        v = self.current_viewer()
        if v:
            v.goto_page(self.page_spin.value() - 1)

    def _goto_current_page(self, page: int):
        v = self.current_viewer()
        if v:
            v.goto_page(page)

    def _on_search_result(self, page: int, rect):
        v = self.current_viewer()
        if v:
            v.goto_page(page)
            self._set_msg(f"Match on page {page + 1}")

    def _set_msg(self, text: str):
        self.status_msg.setText(text)
        log.info(text)

    # ---- view controls -----------------------------------------------------
    def _zoom_in(self):
        v = self.current_viewer()
        if v: v.zoom_in()

    def _zoom_out(self):
        v = self.current_viewer()
        if v: v.zoom_out()

    def _fit_width(self):
        v = self.current_viewer()
        if v: v.fit_width()

    def _fit_page(self):
        v = self.current_viewer()
        if v: v.fit_page()

    def _rotate(self, delta):
        v = self.current_viewer()
        if v:
            v.rotate(delta)
            self.current_doc().dirty = True

    def _set_view_mode(self, mode):
        v = self.current_viewer()
        if v:
            v.set_view_mode(mode)
            self.settings.set("view_mode", mode)

    def _reading_mode(self, mode):
        v = self.current_viewer()
        if v:
            v.set_reading_mode(mode)
        self.settings.set("reading_mode", mode)
        self._set_msg(f"Reading mode: {mode}")

    def _toggle_autoscroll(self, on):
        v = self.current_viewer()
        if v:
            v.toggle_autoscroll(on)

    def _zoom_combo_changed(self):
        text = self.zoom_combo.currentText().strip()
        v = self.current_viewer()
        if not v:
            return
        if text == "Fit Width":
            v.fit_width()
        elif text == "Fit Page":
            v.fit_page()
        else:
            try:
                pct = float(text.replace("%", ""))
                v.set_zoom(pct / 100.0)
            except ValueError:
                pass

    def toggle_fullscreen(self):
        if self.isFullScreen():
            self.showNormal()
        else:
            self.showFullScreen()

    def toggle_theme(self):
        new = "light" if self.settings.get("theme") == "dark" else "dark"
        self.settings.set("theme", new)
        self.apply_theme()

    def apply_theme(self):
        theme = self.settings.get("theme")
        accent = self.settings.get("accent")
        font = self.settings.get("ui_font", "Segoe UI")
        QApplication.instance().setStyleSheet(build_stylesheet(theme, accent, font))

    # ---- page management ---------------------------------------------------
    def insert_blank_page(self):
        doc = self._require_doc()
        if not doc:
            return
        idx = doc.current_page + 1
        w, h = doc.page_size(doc.current_page)
        doc.doc.new_page(pno=idx, width=w, height=h)
        doc.dirty = True
        self._after_page_change()

    def delete_pages(self):
        doc = self._require_doc()
        if not doc:
            return
        text, ok = QInputDialog.getText(
            self, "Delete Pages", "Pages to delete (e.g. 1-3,5):")
        if not ok:
            return
        pages = parse_page_ranges(text, doc.page_count)
        if not pages or len(pages) >= doc.page_count:
            QMessageBox.warning(self, APP_NAME, "Invalid selection (cannot delete all pages).")
            return
        for p in sorted(pages, reverse=True):
            doc.doc.delete_page(p)
        doc.dirty = True
        self._after_page_change()

    def duplicate_page(self):
        doc = self._require_doc()
        if not doc:
            return
        p = doc.current_page
        doc.doc.fullcopy_page(p, p + 1)
        doc.dirty = True
        self._after_page_change()

    def rotate_pages(self):
        doc = self._require_doc()
        if not doc:
            return
        deg, ok = QInputDialog.getItem(
            self, "Rotate Pages", "Rotation:", ["90", "180", "270"], 0, False)
        if not ok:
            return
        text, ok = QInputDialog.getText(self, "Rotate Pages", "Pages (blank = all):")
        if not ok:
            return
        pages = parse_page_ranges(text, doc.page_count)
        for p in pages:
            page = doc.doc[p]
            page.set_rotation((page.rotation + int(deg)) % 360)
        doc.dirty = True
        self._after_page_change()

    def reverse_pages(self):
        doc = self._require_doc()
        if not doc:
            return
        n = doc.page_count
        order = list(range(n - 1, -1, -1))
        try:
            doc.doc.select(order)
            doc.dirty = True
            self._after_page_change()
        except Exception as e:
            QMessageBox.critical(self, "Reverse failed", str(e))

    def _reorder_pages(self, order: List[int]):
        doc = self.current_doc()
        if not doc:
            return
        try:
            doc.doc.select(order)
            doc.dirty = True
            self._after_page_change()
        except Exception as e:
            log.error("reorder: %s", e)

    def extract_pages(self):
        doc = self._require_doc()
        if not doc:
            return
        text, ok = QInputDialog.getText(self, "Extract Pages", "Pages (e.g. 1-3,5):")
        if not ok:
            return
        pages = parse_page_ranges(text, doc.page_count)
        if not pages:
            return
        out, _ = QFileDialog.getSaveFileName(self, "Save extracted PDF", "", "PDF (*.pdf)")
        if not out:
            return
        try:
            new = fitz.open()
            for p in pages:
                new.insert_pdf(doc.doc, from_page=p, to_page=p)
            new.save(out, garbage=4, deflate=True)
            new.close()
            self._set_msg(f"Extracted {len(pages)} pages")
        except Exception as e:
            QMessageBox.critical(self, "Extract failed", str(e))

    def split_pdf(self):
        doc = self._require_doc()
        if not doc:
            return
        every, ok = QInputDialog.getInt(
            self, "Split PDF", "Pages per output file:", 1, 1, doc.page_count)
        if not ok:
            return
        outdir = QFileDialog.getExistingDirectory(self, "Output folder")
        if not outdir:
            return
        try:
            base = os.path.splitext(doc.name)[0]
            part = 1
            for start in range(0, doc.page_count, every):
                end = min(start + every - 1, doc.page_count - 1)
                new = fitz.open()
                new.insert_pdf(doc.doc, from_page=start, to_page=end)
                new.save(os.path.join(outdir, f"{base}_part{part}.pdf"),
                         garbage=4, deflate=True)
                new.close()
                part += 1
            self._set_msg(f"Split into {part - 1} files")
        except Exception as e:
            QMessageBox.critical(self, "Split failed", str(e))

    def merge_pdfs(self):
        doc = self.current_doc()
        paths, _ = QFileDialog.getOpenFileNames(
            self, "Select PDFs to merge (in order)", "", "PDF (*.pdf)")
        if not paths:
            return
        try:
            if doc and doc.doc:
                # Append into the current document.
                for p in paths:
                    with fitz.open(p) as d:
                        doc.doc.insert_pdf(d)
                doc.dirty = True
                self._after_page_change()
                self._set_msg(f"Merged {len(paths)} file(s) into current document")
            else:
                out, _ = QFileDialog.getSaveFileName(self, "Save merged PDF", "", "PDF (*.pdf)")
                if not out:
                    return
                merged = fitz.open()
                for p in paths:
                    with fitz.open(p) as d:
                        merged.insert_pdf(d)
                merged.save(out, garbage=4, deflate=True)
                merged.close()
                self.open_path(out)
        except Exception as e:
            QMessageBox.critical(self, "Merge failed", str(e))

    def _after_page_change(self):
        self._refresh_panels()
        v = self.current_viewer()
        if v:
            v.rebuild()
        self._update_status()

    # ---- insert / annotate -------------------------------------------------
    def insert_image(self):
        doc = self._require_doc()
        if not doc:
            return
        path, _ = QFileDialog.getOpenFileName(
            self, "Insert Image", "", "Images (*.png *.jpg *.jpeg *.bmp)")
        if not path:
            return
        try:
            page = doc.doc[doc.current_page]
            r = page.rect
            w = r.width * 0.5
            h = w * 0.75
            rect = fitz.Rect(r.width * 0.25, r.height * 0.25,
                             r.width * 0.25 + w, r.height * 0.25 + h)
            page.insert_image(rect, filename=path, keep_proportion=True)
            doc.dirty = True
            self._reload_viewer()
            self._set_msg("Image inserted")
        except Exception as e:
            QMessageBox.critical(self, "Insert image failed", str(e))

    def add_text_box(self):
        doc = self._require_doc()
        if not doc:
            return
        text, ok = QInputDialog.getMultiLineText(self, "Add Text", "Text:")
        if not ok or not text:
            return
        size, ok = QInputDialog.getInt(self, "Add Text", "Font size:", 14, 6, 200)
        if not ok:
            return
        try:
            page = doc.doc[doc.current_page]
            r = page.rect
            page.insert_text(fitz.Point(r.width * 0.2, r.height * 0.2),
                             text, fontsize=size, color=(0, 0, 0))
            doc.dirty = True
            self._reload_viewer()
            self._set_msg("Text added")
        except Exception as e:
            QMessageBox.critical(self, "Add text failed", str(e))

    def _quick_annot(self, kind: str):
        """
        Apply a markup annotation to all text on the current page (demonstration).
        In a full editor this would respond to a user-drawn selection.
        """
        doc = self._require_doc()
        if not doc:
            return
        page = doc.doc[doc.current_page]
        try:
            # Annotate the first text line found, as a simple demonstration.
            words = page.get_text("words")
            if not words:
                QMessageBox.information(self, APP_NAME, "No selectable text on this page.")
                return
            rect = fitz.Rect(words[0][:4])
            for w in words[1:8]:
                rect |= fitz.Rect(w[:4])
            if kind == "highlight":
                page.add_highlight_annot(rect)
            elif kind == "underline":
                page.add_underline_annot(rect)
            elif kind == "strikeout":
                page.add_strikeout_annot(rect)
            doc.dirty = True
            self._reload_viewer()
            self.comments_panel.populate(doc)
            self._set_msg(f"{kind.capitalize()} added")
        except Exception as e:
            QMessageBox.critical(self, "Annotation failed", str(e))

    def add_sticky_note(self):
        doc = self._require_doc()
        if not doc:
            return
        text, ok = QInputDialog.getMultiLineText(self, "Sticky Note", "Note text:")
        if not ok:
            return
        try:
            page = doc.doc[doc.current_page]
            page.add_text_annot(fitz.Point(40, 40), text)
            doc.dirty = True
            self._reload_viewer()
            self.comments_panel.populate(doc)
            self._set_msg("Sticky note added")
        except Exception as e:
            QMessageBox.critical(self, "Note failed", str(e))

    def _reload_viewer(self):
        v = self.current_viewer()
        if v:
            for lbl in v._labels:
                lbl.rendered_zoom = -1.0
            v._render_visible()

    # ---- tools -------------------------------------------------------------
    def show_watermark(self):
        doc = self._require_doc()
        if not doc:
            return
        dlg = WatermarkDialog(doc, self)
        if dlg.exec() == QDialog.Accepted:
            dlg.apply()
            self._reload_viewer()
            self._set_msg("Watermark applied")

    def show_header_footer(self):
        doc = self._require_doc()
        if not doc:
            return
        dlg = HeaderFooterDialog(doc, self)
        if dlg.exec() == QDialog.Accepted:
            dlg.apply()
            self._reload_viewer()
            self._set_msg("Header/Footer applied")

    def show_compress(self):
        doc = self._require_doc()
        if not doc:
            return
        CompressDialog(doc, self).exec()

    def show_ocr(self):
        doc = self._require_doc()
        if not doc:
            return
        OcrDialog(doc, self).exec()

    def show_batch(self):
        BatchDialog(self).exec()

    def show_security(self):
        doc = self._require_doc()
        if not doc:
            return
        dlg = SecurityDialog(doc, self)
        if dlg.exec() != QDialog.Accepted:
            return
        path, _ = QFileDialog.getSaveFileName(
            self, "Save secured PDF", doc.path or "secured.pdf", "PDF (*.pdf)")
        if not path:
            return
        try:
            kwargs = dlg.build_save_kwargs()
            doc.doc.save(path, garbage=4, deflate=True, **kwargs)
            self._set_msg("Security settings applied")
            QMessageBox.information(self, "Security", f"Saved to {os.path.basename(path)}")
        except Exception as e:
            QMessageBox.critical(self, "Security failed", str(e))

    def show_settings(self):
        dlg = SettingsDialog(self.settings, self)
        if dlg.exec() == QDialog.Accepted:
            self.apply_theme()
            self._reconfigure_autosave()
            self._set_msg("Settings saved")

    def show_about(self):
        deps = []
        for name, ok in [("PyMuPDF", HAS_FITZ), ("pypdf", HAS_PYPDF),
                         ("Pillow", HAS_PIL), ("ReportLab", HAS_REPORTLAB),
                         ("pytesseract", HAS_TESSERACT), ("OpenCV", HAS_CV2),
                         ("python-docx", HAS_DOCX), ("openpyxl", HAS_OPENPYXL),
                         ("python-pptx", HAS_PPTX)]:
            deps.append(f"  {'✓' if ok else '✗'} {name}")
        QMessageBox.about(
            self, f"About {APP_NAME}",
            f"<h3>{APP_NAME} {APP_VERSION}</h3>"
            "<p>A professional desktop PDF application built with PySide6.</p>"
            "<p><b>Optional components:</b></p>"
            "<pre>" + "\n".join(deps) + "</pre>"
            f"<p style='color:gray'>Logs: {LOG_PATH}</p>")

    # ---- converters --------------------------------------------------------
    def _convert(self, kind: str):
        doc = self._require_doc()
        if not doc:
            return
        try:
            if kind == "text":
                out, _ = QFileDialog.getSaveFileName(self, "Save Text", "", "Text (*.txt)")
                if out: self._run_worker(pdf_to_text, doc, out, label="Converting to text")
            elif kind == "html":
                out, _ = QFileDialog.getSaveFileName(self, "Save HTML", "", "HTML (*.html)")
                if out: self._run_worker(pdf_to_html, doc, out, label="Converting to HTML")
            elif kind == "image":
                outdir = QFileDialog.getExistingDirectory(self, "Output folder")
                if outdir: self._run_worker(pdf_to_images, doc, outdir, label="Rendering images")
            elif kind == "word":
                out, _ = QFileDialog.getSaveFileName(self, "Save Word", "", "Word (*.docx)")
                if out: self._run_worker(pdf_to_word, doc, out, label="Converting to Word")
            elif kind == "excel":
                out, _ = QFileDialog.getSaveFileName(self, "Save Excel", "", "Excel (*.xlsx)")
                if out: self._run_worker(pdf_to_excel, doc, out, label="Converting to Excel")
            elif kind == "ppt":
                out, _ = QFileDialog.getSaveFileName(self, "Save PowerPoint", "", "PowerPoint (*.pptx)")
                if out: self._run_worker(pdf_to_powerpoint, doc, out, label="Converting to PowerPoint")
        except Exception as e:
            QMessageBox.critical(self, "Conversion failed", str(e))

    def images_to_pdf_action(self):
        if not self._ensure_fitz():
            return
        paths, _ = QFileDialog.getOpenFileNames(
            self, "Select images", "", "Images (*.png *.jpg *.jpeg *.bmp *.tiff)")
        if not paths:
            return
        out, _ = QFileDialog.getSaveFileName(self, "Save PDF", "", "PDF (*.pdf)")
        if not out:
            return
        self._run_worker(images_to_pdf, paths, out, label="Building PDF",
                         on_done=lambda r: self.open_path(out))

    def text_to_pdf_action(self):
        if not self._ensure_fitz():
            return
        path, _ = QFileDialog.getOpenFileName(self, "Select text file", "", "Text (*.txt)")
        if not path:
            return
        out, _ = QFileDialog.getSaveFileName(self, "Save PDF", "", "PDF (*.pdf)")
        if not out:
            return
        self._run_worker(text_to_pdf, path, out, label="Building PDF",
                         on_done=lambda r: self.open_path(out))

    # ---- worker plumbing ---------------------------------------------------
    def _run_worker(self, fn, *args, label="Working…", on_done=None):
        progress = QProgressDialog(label, "Cancel", 0, 100, self)
        progress.setWindowModality(Qt.WindowModal)
        progress.setMinimumDuration(0)
        progress.setValue(0)

        worker = Worker(fn, *args)

        def _on_progress(pct, msg):
            progress.setValue(pct)
            if msg:
                progress.setLabelText(f"{label}\n{msg}")

        def _on_finished(result):
            progress.setValue(100)
            progress.close()
            self._set_msg(f"{label} — done")
            if on_done:
                on_done(result)
            else:
                QMessageBox.information(self, APP_NAME, f"{label} complete.")

        def _on_error(msg):
            progress.close()
            QMessageBox.critical(self, "Operation failed", msg)

        worker.signals.progress.connect(_on_progress)
        worker.signals.finished.connect(_on_finished)
        worker.signals.error.connect(_on_error)
        progress.canceled.connect(lambda: None)  # workers are best-effort cancellable
        self.threadpool.start(worker)

    # ---- autosave ----------------------------------------------------------
    def _reconfigure_autosave(self):
        if self.settings.get_bool("autosave", True):
            interval = self.settings.get_int("autosave_interval", 180) * 1000
            self.autosave_timer.start(max(30000, interval))
        else:
            self.autosave_timer.stop()

    def _autosave(self):
        recovery_dir = os.path.join(APP_DIR, "autosave")
        os.makedirs(recovery_dir, exist_ok=True)
        for i in range(self.tabs.count()):
            w = self.tabs.widget(i)
            if isinstance(w, DocumentTab) and w.document.dirty and w.document.doc:
                try:
                    name = w.document.name.replace(".pdf", "") + "_autosave.pdf"
                    w.document.doc.save(os.path.join(recovery_dir, name),
                                        garbage=3, deflate=True)
                except Exception as e:
                    log.error("autosave: %s", e)
        self._set_msg("Autosaved")

    # ---- misc --------------------------------------------------------------
    def _not_impl(self):
        self._set_msg("That action is not yet implemented.")

    def _ensure_fitz(self) -> bool:
        if not HAS_FITZ:
            QMessageBox.critical(
                self, APP_NAME,
                "PyMuPDF (fitz) is required for this operation.\n"
                "Install it with:  pip install PyMuPDF")
            return False
        return True

    # ---- drag & drop -------------------------------------------------------
    def dragEnterEvent(self, event: QDragEnterEvent):
        if event.mimeData().hasUrls():
            event.acceptProposedAction()

    def dropEvent(self, event: QDropEvent):
        for url in event.mimeData().urls():
            path = url.toLocalFile()
            if path.lower().endswith(".pdf"):
                self.open_path(path)
            elif path.lower().endswith((".png", ".jpg", ".jpeg", ".bmp")):
                # Offer to convert dropped image to PDF
                out = os.path.splitext(path)[0] + ".pdf"
                try:
                    images_to_pdf([path], out)
                    self.open_path(out)
                except Exception as e:
                    QMessageBox.critical(self, "Image import failed", str(e))

    def closeEvent(self, event):
        # Prompt for unsaved tabs.
        for i in range(self.tabs.count()):
            w = self.tabs.widget(i)
            if isinstance(w, DocumentTab) and w.document.dirty:
                ret = QMessageBox.question(
                    self, "Unsaved changes",
                    "There are unsaved documents. Quit anyway?",
                    QMessageBox.Yes | QMessageBox.No)
                if ret == QMessageBox.No:
                    event.ignore()
                    return
                break
        event.accept()


# ----------------------------------------------------------------------------
# 11. Entry point
# ----------------------------------------------------------------------------
def main():
    QApplication.setApplicationName(APP_NAME)
    QApplication.setOrganizationName(ORG_NAME)
    QApplication.setApplicationVersion(APP_VERSION)
    # High-DPI friendliness
    if hasattr(Qt, "AA_EnableHighDpiScaling"):
        QApplication.setAttribute(Qt.AA_EnableHighDpiScaling, True)
    if hasattr(Qt, "AA_UseHighDpiPixmaps"):
        QApplication.setAttribute(Qt.AA_UseHighDpiPixmaps, True)

    app = QApplication(sys.argv)
    app.setStyle(QStyleFactory.create("Fusion"))

    settings = SettingsManager(os.path.join(APP_DIR, "settings.db"))

    # Record exactly which interpreter is running so dependency issues are
    # diagnosable (multiple Python installs can shadow each other on Windows).
    log.info("Interpreter: %s", sys.executable)
    log.info("Version: %s", sys.version.replace("\n", " "))
    log.info("HAS_FITZ=%s  sys.prefix=%s", HAS_FITZ, sys.prefix)

    window = MainWindow(settings)
    install_excepthook(lambda: window)
    window.apply_theme()
    window.show()

    if not HAS_FITZ:
        QMessageBox.warning(
            window, APP_NAME,
            "PyMuPDF (fitz) is not installed for THIS Python interpreter — "
            "most features are disabled.\n\n"
            f"Running interpreter:\n{sys.executable}\n"
            f"Python {sys.version.split()[0]}\n\n"
            "Install into THIS interpreter with:\n"
            f'"{sys.executable}" -m pip install PySide6 PyMuPDF pypdf Pillow '
            "reportlab pytesseract opencv-python python-docx openpyxl python-pptx")

    # Open any files passed on the command line.
    for arg in sys.argv[1:]:
        if arg.lower().endswith(".pdf") and os.path.exists(arg):
            window.open_path(arg)

    sys.exit(app.exec())


if __name__ == "__main__":
    main()
